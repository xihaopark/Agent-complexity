from __future__ import annotations

import json
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def read_json(path: Path) -> Dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def safe_float(value: Any, default: float = 0.0) -> float:
    try:
        if value is None or value == "":
            return default
        return float(value)
    except Exception:
        return default


def safe_int(value: Any, default: int = 0) -> int:
    try:
        if value is None or value == "":
            return default
        return int(value)
    except Exception:
        return default


def mmss(seconds: float) -> str:
    if seconds <= 0:
        return "-"
    m = int(seconds // 60)
    s = int(seconds - m * 60)
    return f"{m}:{s:02d}"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)

    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(12.0), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 247, 250)
    box.line.color.rgb = RGBColor(220, 225, 232)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(15, 23, 42)

    sub = slide.shapes.add_textbox(Inches(0.9), Inches(3.1), Inches(12.0), Inches(0.8))
    st = sub.text_frame
    st.clear()
    p2 = st.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(55, 65, 81)


def add_section_header(slide, title: str) -> None:
    header = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(12.2), Inches(0.6))
    tf = header.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(15, 23, 42)


def add_bullets(slide, left: float, top: float, width: float, height: float, bullets: List[str]) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(31, 41, 55)


def add_table(
    slide,
    *,
    left: float,
    top: float,
    col_widths: List[float],
    row_heights: List[float],
    data: List[List[str]],
    header_rows: int = 1,
    cell_fill: Dict[Tuple[int, int], RGBColor] | None = None,
    cell_font_color: Dict[Tuple[int, int], RGBColor] | None = None,
) -> None:
    rows = len(data)
    cols = len(data[0]) if rows else 0
    width = sum(col_widths)
    height = sum(row_heights)
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)
    for i, h in enumerate(row_heights):
        table.rows[i].height = Inches(h)

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(data[i][j])
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if i < header_rows or j > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(12 if i >= header_rows else 12)
                    run.font.bold = True if i < header_rows else False
                    run.font.color.rgb = RGBColor(15, 23, 42) if i < header_rows else RGBColor(31, 41, 55)

            if i < header_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            if cell_fill and (i, j) in cell_fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_fill[(i, j)]
            if cell_font_color and (i, j) in cell_font_color:
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.color.rgb = cell_font_color[(i, j)]


def build_pivots(rows: List[Dict[str, Any]]) -> Dict[str, Any]:
    workflows = sorted({r["workflow_id"] for r in rows})
    agents = sorted({r["agent_name"] for r in rows})
    by_key: Dict[Tuple[str, str], Dict[str, Any]] = {(r["workflow_id"], r["agent_name"]): r for r in rows}

    def get(workflow: str, agent: str, key: str, default: Any = "") -> Any:
        r = by_key.get((workflow, agent)) or {}
        return r.get(key, default)

    status = [[get(w, a, "status", "") for a in agents] for w in workflows]
    duration = [[safe_float(get(w, a, "duration_seconds", 0.0)) for a in agents] for w in workflows]
    llm_calls = [[safe_int(get(w, a, "llm_call_count", 0)) for a in agents] for w in workflows]
    tokens = [[safe_int(get(w, a, "total_tokens", 0)) for a in agents] for w in workflows]
    cost = [[safe_float(get(w, a, "cost", 0.0)) for a in agents] for w in workflows]

    return {
        "workflows": workflows,
        "agents": agents,
        "status": status,
        "duration": duration,
        "llm_calls": llm_calls,
        "tokens": tokens,
        "cost": cost,
    }


def extract_nospace_reason(agent_run_path: Path) -> str:
    try:
        text = agent_run_path.read_text(encoding="utf-8")
    except Exception:
        return ""
    m = re.search(r"NoSpaceLeftError: No space left on devices\\.", text)
    if not m:
        return ""
    start = max(0, m.start() - 220)
    end = min(len(text), m.end() + 220)
    snippet = text[start:end]
    snippet = snippet.replace("\\n", "\n")
    snippet = re.sub(r"[ \\t]+", " ", snippet)
    return snippet.strip()


def main() -> int:
    input_root = Path("finish/Renzo_DA_Agent/data/formal_peer_runs/20260408_formal").resolve()
    summary_flat = input_root / "summary_flat.json"
    rows = read_json(summary_flat)
    pivots = build_pivots(rows)
    workflows: List[str] = pivots["workflows"]
    agents: List[str] = pivots["agents"]

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        title="Formal Workflow Runs: Renzo vs Peer Agents",
        subtitle=f"Run root: {input_root.name} • Generated: {now_iso()}",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Experiment Setup")
    add_bullets(
        slide,
        left=0.9,
        top=1.3,
        width=12.0,
        height=5.5,
        bullets=[
            f"Agents: {', '.join(agents)}",
            f"Workflows: {', '.join(workflows)}",
            "Mode: agent-only (no direct baseline), per-workflow isolated workdir",
            "Framework enforcement: biomni/stella require real framework code path",
            "Metrics: status, step completion, duration (process wall time), LLM call count, token/cost (when available)",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Outcome Matrix (Success/Fail)")
    header = ["workflow_id", *agents]
    data = [header]
    fills: Dict[Tuple[int, int], RGBColor] = {}
    fcolors: Dict[Tuple[int, int], RGBColor] = {}
    for i, w in enumerate(workflows):
        row = [w]
        for j, a in enumerate(agents):
            s = pivots["status"][i][j]
            row.append(s)
        data.append(row)
    for i in range(1, len(data)):
        for j in range(1, len(data[0])):
            val = data[i][j].lower()
            if val == "success":
                fills[(i, j)] = RGBColor(220, 252, 231)
                fcolors[(i, j)] = RGBColor(22, 101, 52)
            elif val == "failed":
                fills[(i, j)] = RGBColor(254, 226, 226)
                fcolors[(i, j)] = RGBColor(153, 27, 27)
            else:
                fills[(i, j)] = RGBColor(243, 244, 246)
                fcolors[(i, j)] = RGBColor(55, 65, 81)
    add_table(
        slide,
        left=0.8,
        top=1.3,
        col_widths=[3.6] + [1.6] * len(agents),
        row_heights=[0.45] + [0.42] * len(workflows),
        data=data,
        cell_fill=fills,
        cell_font_color=fcolors,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Runtime (Duration per Workflow)")
    header = ["workflow_id", *agents]
    data = [header]
    for i, w in enumerate(workflows):
        row = [w]
        for j, _a in enumerate(agents):
            row.append(mmss(pivots["duration"][i][j]))
        data.append(row)
    add_table(
        slide,
        left=0.8,
        top=1.3,
        col_widths=[3.6] + [1.6] * len(agents),
        row_heights=[0.45] + [0.42] * len(workflows),
        data=data,
    )
    add_bullets(
        slide,
        left=0.9,
        top=6.35,
        width=12.0,
        height=0.9,
        bullets=["Note: durations include workflow engine + environment provisioning (e.g., conda env creation)."],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "LLM Usage (Calls / Tokens / Cost)")
    header = ["workflow_id", "agent", "status", "llm_calls", "total_tokens", "cost($)"]
    data = [header]
    for w in workflows:
        for a in agents:
            r = next((x for x in rows if x["workflow_id"] == w and x["agent_name"] == a), None) or {}
            data.append(
                [
                    w,
                    a,
                    str(r.get("status") or ""),
                    str(r.get("llm_call_count") or 0),
                    str(r.get("total_tokens") or 0),
                    f"{safe_float(r.get('cost', 0.0)):.6f}",
                ]
            )
    add_table(
        slide,
        left=0.8,
        top=1.3,
        col_widths=[3.3, 1.4, 1.2, 1.2, 1.6, 1.4],
        row_heights=[0.45] + [0.28] * min(18, len(data) - 1),
        data=data[:19],
        header_rows=1,
    )
    add_bullets(
        slide,
        left=0.9,
        top=6.35,
        width=12.0,
        height=0.9,
        bullets=[
            "Note: token/cost are 0 for some runs because framework-native calls do not always surface usage metadata.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Why did Snakemake Workflow Template fail for all agents?")
    templ_dir = input_root / "snakemake-workflow-template-finish" / "artifacts"
    renzo_path = templ_dir / "agent-renzo-run.json"
    snippet = extract_nospace_reason(renzo_path)
    bullets = [
        "Root cause: Snakemake conda environment provisioning failed with NoSpaceLeftError during prepare_reference.",
        "This is an infra/resource failure (disk space), not an agent-specific planning error.",
        "Typical fix: clear runtime cache, move conda-prefix to larger disk, or pre-build envs/images.",
    ]
    add_bullets(slide, left=0.9, top=1.3, width=12.2, height=2.2, bullets=bullets)
    if snippet:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.7), Inches(12.1), Inches(3.3))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(15, 23, 42)
        box.line.color.rgb = RGBColor(15, 23, 42)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = snippet[:900]
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(226, 232, 240)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Related Failure: zarp-finish")
    zarp_dir = input_root / "zarp-finish" / "artifacts"
    zarp_renzo_path = zarp_dir / "agent-renzo-run.json"
    zarp_snippet = extract_nospace_reason(zarp_renzo_path)
    add_bullets(
        slide,
        left=0.9,
        top=1.3,
        width=12.2,
        height=2.0,
        bullets=[
            "Observed same pattern: conda env provisioning failed with NoSpaceLeftError (e.g., cutadapt.yaml).",
            "Impact: all agents fail early because step environments cannot be created.",
            "Fix is identical: free disk / relocate cache / pre-build envs.",
        ],
    )
    if zarp_snippet:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.7), Inches(12.1), Inches(3.3))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(15, 23, 42)
        box.line.color.rgb = RGBColor(15, 23, 42)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = zarp_snippet[:900]
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(226, 232, 240)

    out_dir = Path("finish/Renzo_DA_Agent/data/reports").resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"formal_peer_runs_{input_root.name}.pptx"
    prs.save(str(out_path))
    print(str(out_path))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
