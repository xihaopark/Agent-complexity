from __future__ import annotations

import json
import re
import importlib.util
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Tuple

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
INPUT_PPTX = ROOT / "ppt_assets" / "outputs" / "academic-agent-fairness-v3-results.pptx"
RUN_ROOT = ROOT / "finish" / "Renzo_DA_Agent" / "data" / "formal_peer_runs" / "20260408_formal"
SUMMARY_FLAT = RUN_ROOT / "summary_flat.json"
V3_SCRIPT = ROOT / "ppt_assets" / "workspace" / "generate_academic_agent_ppt_v3_results.py"


def load_v3() -> Any:
    spec = importlib.util.spec_from_file_location("academic_v3", str(V3_SCRIPT))
    if spec is None or spec.loader is None:
        raise RuntimeError(f"cannot load v3 script: {V3_SCRIPT}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def mmss(seconds: float | None) -> str:
    if seconds is None:
        return "-"
    try:
        s = float(seconds)
    except Exception:
        return "-"
    if s <= 0:
        return "-"
    m = int(s // 60)
    sec = int(s - m * 60)
    return f"{m}:{sec:02d}"


def build_matrix(rows: List[Dict[str, Any]]) -> Dict[str, Any]:
    workflows = sorted({r["workflow_id"] for r in rows})
    agents = sorted({r["agent_name"] for r in rows})
    by_key: Dict[Tuple[str, str], Dict[str, Any]] = {(r["workflow_id"], r["agent_name"]): r for r in rows}

    def get(workflow: str, agent: str) -> Dict[str, Any]:
        return by_key.get((workflow, agent)) or {}

    return {"workflows": workflows, "agents": agents, "get": get}


def first_run(paragraph):
    return paragraph.runs[0] if paragraph.runs else paragraph.add_run()


def extract_nospace_snippet(path: Path) -> str:
    if not path.exists():
        return ""
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""
    m = re.search(r"NoSpaceLeftError: No space left on devices\\.", text)
    if not m:
        return ""
    start = max(0, m.start() - 220)
    end = min(len(text), m.end() + 220)
    snippet = text[start:end]
    snippet = snippet.replace("\r\n", "\n").replace("\r", "\n")
    snippet = re.sub(r"[ \\t]+", " ", snippet)
    return snippet.strip()


def add_peer_results_overview(prs: Presentation, page: int, rows: List[Dict[str, Any]]) -> int:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    v3.set_bg(slide)
    v3.add_title(
        slide,
        "新增实验：Renzo vs 4 个 Peer Agents（Agent-only）",
        f"Run root: {RUN_ROOT.name} • 生成于 {now_iso()}",
    )

    matrix = build_matrix(rows)
    workflows = matrix["workflows"]
    agents = matrix["agents"]

    table = slide.shapes.add_table(
        1 + len(workflows),
        1 + len(agents),
        Inches(0.42),
        Inches(1.38),
        Inches(12.50),
        Inches(2.85),
    ).table

    table.columns[0].width = Inches(3.2)
    for i in range(len(agents)):
        table.columns[i + 1].width = Inches((12.50 - 3.2) / max(1, len(agents)))

    headers = ["Workflow", *agents]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = v3.BLUE_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = first_run(p)
        r.font.name = v3.FONT
        r.font.bold = True
        r.font.size = Pt(10.2)
        r.font.color.rgb = v3.WHITE

    for r_idx, wid in enumerate(workflows, start=1):
        table.cell(r_idx, 0).text = v3.short_workflow_name(wid)
        for c in range(0, len(agents) + 1):
            cell = table.cell(r_idx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = v3.WHITE if r_idx % 2 else v3.SOFT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = first_run(p)
            run.font.name = v3.FONT
            run.font.size = Pt(9.0)
            run.font.color.rgb = v3.TEXT

        for c_idx, agent in enumerate(agents, start=1):
            r = matrix["get"](wid, agent)
            status = str(r.get("status") or "")
            cell = table.cell(r_idx, c_idx)
            cell.text = status
            p = cell.text_frame.paragraphs[0]
            run = first_run(p)
            run.font.color.rgb = v3.GREEN if status == "success" else v3.RED

    v3.add_card(
        slide,
        Inches(0.62),
        Inches(4.45),
        Inches(4.0),
        Inches(1.35),
        "要点",
        "Peer agents 与 Renzo 并排运行 workflow；执行引擎一致，差异主要来自编排/选步策略与框架行为。",
        accent=v3.BLUE,
        body_size=10.6,
    )
    v3.add_card(
        slide,
        Inches(4.72),
        Inches(4.45),
        Inches(4.0),
        Inches(1.35),
        "失败集中",
        "snakemake-workflow-template 与 zarp 的失败属于环境/磁盘资源问题（conda env 创建失败），与 agent 策略无关。",
        accent=v3.RED,
        body_size=10.4,
    )
    v3.add_card(
        slide,
        Inches(8.82),
        Inches(4.45),
        Inches(4.0),
        Inches(1.35),
        "指标说明",
        "token/cost 对 Biomni/STELLA 可能缺失（真实框架未回传 usage）；latency 采用 process wall-time。",
        accent=v3.PURPLE,
        body_size=10.3,
    )

    v3.add_footer(slide, page)
    return page + 1


def add_peer_metrics_slide(prs: Presentation, page: int, rows: List[Dict[str, Any]]) -> int:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    v3.set_bg(slide)
    v3.add_title(slide, "Peer Agents 指标：时长 / LLM 调用次数 / Tokens", "单位：时长=mm:ss；tokens/cost 仅在可用时统计")

    matrix = build_matrix(rows)
    workflows = matrix["workflows"]
    agents = matrix["agents"]

    headers = ["Workflow", "Agent", "Status", "Time", "LLM calls", "Tokens", "Cost"]
    max_rows = min(16, len(workflows) * len(agents))
    table = slide.shapes.add_table(
        1 + max_rows,
        len(headers),
        Inches(0.42),
        Inches(1.38),
        Inches(12.50),
        Inches(5.35),
    ).table
    widths = [2.8, 1.3, 0.9, 0.9, 1.1, 1.1, 1.1]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = v3.BLUE_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = first_run(p)
        r.font.name = v3.FONT
        r.font.bold = True
        r.font.size = Pt(10.2)
        r.font.color.rgb = v3.WHITE

    out_rows: List[List[str]] = []
    for wid in workflows:
        for agent in agents:
            r = matrix["get"](wid, agent)
            out_rows.append(
                [
                    v3.short_workflow_name(wid),
                    agent,
                    str(r.get("status") or ""),
                    mmss(r.get("duration_seconds")),
                    str(r.get("llm_call_count") or 0),
                    str(r.get("total_tokens") or 0),
                    f"${float(r.get('cost') or 0.0):.4f}",
                ]
            )

    for i, vals in enumerate(out_rows[:max_rows], start=1):
        for c, text in enumerate(vals):
            cell = table.cell(i, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = v3.WHITE if i % 2 else v3.SOFT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c in {0, 1} else PP_ALIGN.CENTER
            r = first_run(p)
            r.font.name = v3.FONT
            r.font.size = Pt(9.0)
            if c == 2:
                r.font.color.rgb = v3.GREEN if text == "success" else v3.RED
            else:
                r.font.color.rgb = v3.TEXT

    v3.add_footer(slide, page)
    return page + 1


def add_failure_rootcause_slide(prs: Presentation, page: int) -> int:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    v3.set_bg(slide)
    v3.add_title(slide, "根因：为什么 Template / Zarp 会全失败？", "共因：conda env provisioning 时磁盘空间耗尽 (NoSpaceLeftError)")

    tmpl = RUN_ROOT / "snakemake-workflow-template-finish" / "artifacts" / "agent-renzo-run.json"
    zarp = RUN_ROOT / "zarp-finish" / "artifacts" / "agent-renzo-run.json"
    tmpl_snip = extract_nospace_snippet(tmpl)
    zarp_snip = extract_nospace_snippet(zarp)

    v3.add_card(
        slide,
        Inches(0.62),
        Inches(1.55),
        Inches(12.20),
        Inches(1.15),
        "解释",
        "失败发生在 Snakemake 创建 conda env（基于 envs/*.yaml）阶段；由于缓存/conda 前缀占满磁盘，导致 token 文件写入失败。",
        accent=v3.RED,
        body_size=11.0,
    )

    if tmpl_snip:
        v3.add_textbox(slide, Inches(0.62), Inches(2.88), Inches(12.20), Inches(1.15), "Template 日志片段：", size=11, bold=True, color=v3.BLUE_DARK)
        v3.add_textbox(slide, Inches(0.62), Inches(3.15), Inches(12.20), Inches(1.15), tmpl_snip[:260], size=9.2, color=v3.MUTED, align=PP_ALIGN.LEFT)
    if zarp_snip:
        v3.add_textbox(slide, Inches(0.62), Inches(4.38), Inches(12.20), Inches(0.32), "Zarp 日志片段：", size=11, bold=True, color=v3.BLUE_DARK)
        v3.add_textbox(slide, Inches(0.62), Inches(4.65), Inches(12.20), Inches(1.15), zarp_snip[:260], size=9.2, color=v3.MUTED, align=PP_ALIGN.LEFT)

    v3.add_card(
        slide,
        Inches(0.62),
        Inches(6.02),
        Inches(12.20),
        Inches(0.95),
        "建议",
        "清理 runtime cache；并使用 shared conda prefix（避免每个 agent 复制一套 env）。必要时把 conda-prefix 移到更大磁盘。",
        accent=v3.GREEN,
        body_size=11.0,
    )
    v3.add_footer(slide, page)
    return page + 1


def main() -> None:
    if not INPUT_PPTX.exists():
        raise SystemExit(f"missing pptx: {INPUT_PPTX}")
    if not SUMMARY_FLAT.exists():
        raise SystemExit(f"missing summary: {SUMMARY_FLAT}")
    global v3
    v3 = load_v3()

    prs = Presentation(str(INPUT_PPTX))
    rows = read_json(SUMMARY_FLAT)
    page = len(prs.slides) + 1

    page = add_peer_results_overview(prs, page, rows)
    page = add_peer_metrics_slide(prs, page, rows)
    page = add_failure_rootcause_slide(prs, page)

    prs.save(str(INPUT_PPTX))
    print(f"updated {INPUT_PPTX}")


if __name__ == "__main__":
    main()
