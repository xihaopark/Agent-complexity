from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUT = Path("/lab_workspace/projects/Agent-complexity/main/ppt_assets/outputs/academic-agent-fairness-v2.pptx")

W = Inches(13.333)
H = Inches(7.5)

BLUE = RGBColor(37, 99, 235)
BLUE_DARK = RGBColor(30, 64, 175)
BLUE_LIGHT = RGBColor(219, 234, 254)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
GRAY = RGBColor(226, 232, 240)
GRAY_DARK = RGBColor(148, 163, 184)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(220, 38, 38)
AMBER = RGBColor(217, 119, 6)
GREEN = RGBColor(22, 163, 74)

FONT = "Noto Sans CJK SC"


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=20, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullet_box(slide, left, top, width, height, bullets, size=18, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.bullet = True
        p.space_after = Pt(7)
    return box


def add_rect(slide, left, top, width, height, fill, line=GRAY, radius=None):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def style_text_in_shape(shape, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = shape.text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_shape_text(shape, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.6), title, size=26, bold=True)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(0.98), Inches(12.1), Inches(0.03)).fill.solid()
    line = slide.shapes[-1]
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.05), Inches(10.8), Inches(0.35), subtitle, size=12, color=MUTED)


def add_footer(slide, page, note="Agent Benchmarking for Bioinformatics Workflows"):
    add_textbox(slide, Inches(0.6), Inches(7.08), Inches(7.5), Inches(0.2), note, size=9, color=GRAY_DARK)
    add_textbox(slide, Inches(12.1), Inches(7.02), Inches(0.5), Inches(0.25), str(page), size=11, color=GRAY_DARK, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, body, accent=BLUE):
    shape = add_rect(slide, left, top, width, height, WHITE, line=GRAY, radius=True)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.08), height).fill.solid()
    bar = slide.shapes[-1]
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_textbox(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.3), Inches(0.28), title, size=13, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.48), width - Inches(0.3), height - Inches(0.55), body, size=11, color=MUTED)
    return shape


def add_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Toward Fair Evaluation of AI Agents on Bioinformatics Workflows",
              "为什么需要从“能否完成任务”转向“能否在受控预算下公平地完成真实 workflow”")

    add_textbox(slide, Inches(0.7), Inches(1.35), Inches(6.3), Inches(0.7),
                "研究动机", size=22, bold=True, color=BLUE_DARK)
    add_bullet_box(
        slide, Inches(0.8), Inches(1.95), Inches(5.8), Inches(2.2),
        [
            "AI agent 已从单轮问答扩展到多步规划、工具调用与环境交互。",
            "生物信息学 workflow 具有长链条、强依赖、文件密集和可复现性敏感等特点。",
            "现有 benchmark 很少同时控制同一 workflow、同一工具权限、同一预算与同一观测口径。"
        ],
        size=17
    )

    pipe_y = Inches(1.55)
    boxes = [
        ("研究问题", BLUE_LIGHT),
        ("Agent 规划", BLUE_LIGHT),
        ("Workflow 执行", BLUE_LIGHT),
        ("结果评估", BLUE_LIGHT),
    ]
    x = Inches(7.0)
    for idx, (label, color) in enumerate(boxes):
        box = add_rect(slide, x + Inches(idx * 1.3), pipe_y, Inches(1.12), Inches(0.68), color, line=BLUE, radius=True)
        set_shape_text(box, label, size=12, bold=True, color=TEXT)
        if idx < len(boxes) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + Inches(idx * 1.3 + 1.05), pipe_y + Inches(0.2), Inches(0.25), Inches(0.28))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BLUE
            arrow.line.color.rgb = BLUE

    add_card(slide, Inches(6.95), Inches(2.65), Inches(1.8), Inches(1.2),
             "GAIA", "466 real-world questions\nHuman 92%\nGPT-4+plugins 15%", accent=BLUE)
    add_card(slide, Inches(8.95), Inches(2.65), Inches(1.8), Inches(1.2),
             "OSWorld", "369 desktop tasks\nHuman 72.36%\nBest model 12.24%", accent=GREEN)
    add_card(slide, Inches(10.95), Inches(2.65), Inches(1.8), Inches(1.2),
             "BixBench", "53 bioinformatics scenarios\n296 questions\nfocus on multi-step analysis", accent=AMBER)

    quote = add_rect(slide, Inches(0.8), Inches(4.55), Inches(12.0), Inches(1.4), BLUE_DARK, line=BLUE_DARK, radius=True)
    set_shape_text(
        quote,
        "核心问题：如果不同 agent 在“同一 workflow”上使用不同环境、不同观测方式和不同预算，最终比较结果往往并不公平。",
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )

    add_footer(slide, 1)


def add_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Benchmark Landscape: What Existing Benchmarks Cover, and What They Miss")

    rows = 6
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(0.65), Inches(1.35), Inches(7.1), Inches(4.4)).table
    headers = ["Benchmark", "主要能力", "优点", "对本课题的缺口"]
    data = [
        ["GAIA", "多步推理 + 工具", "真实世界问题", "非 workflow 场景，环境控制弱"],
        ["SWE-bench", "仓库级修复", "长链工程任务", "偏软件工程，不含数据流水线"],
        ["WebArena", "Web 交互", "可复现网页环境", "缺少文件密集型科学 workflow"],
        ["OSWorld", "桌面/多应用", "真实计算机环境", "领域约束弱，生信工具链缺失"],
        ["BixBench", "生信分析", "贴近 bioinformatics", "尚未强调公平的 agent-to-agent 协议"],
    ]
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
    for c in range(cols):
        p = table.cell(0, c).text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = FONT
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(248, 250, 252)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.name = FONT
            run.font.size = Pt(10.5)
            run.font.color.rgb = TEXT

    add_textbox(slide, Inches(8.15), Inches(1.42), Inches(4.1), Inches(0.35), "真实性 vs 领域适配度", size=16, bold=True)
    # quadrant
    left = Inches(8.2)
    top = Inches(1.9)
    width = Inches(4.0)
    height = Inches(3.0)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height).fill.background()
    frame = slide.shapes[-1]
    frame.line.color.rgb = GRAY_DARK
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + width / 2, top, Inches(0.01), height).fill.solid()
    vline = slide.shapes[-1]
    vline.fill.fore_color.rgb = GRAY
    vline.line.color.rgb = GRAY
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top + height / 2, width, Inches(0.01)).fill.solid()
    hline = slide.shapes[-1]
    hline.fill.fore_color.rgb = GRAY
    hline.line.color.rgb = GRAY
    add_textbox(slide, left + Inches(1.25), top + Inches(3.02), Inches(1.6), Inches(0.2), "真实性", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, left - Inches(0.1), top + Inches(1.3), Inches(0.35), Inches(0.7), "领域适配度", size=10, color=MUTED)

    points = [
        ("GAIA", 0.9, 1.0, BLUE),
        ("SWE-bench", 1.8, 1.5, GREEN),
        ("WebArena", 2.6, 1.1, AMBER),
        ("OSWorld", 3.0, 1.5, RED),
        ("BixBench", 3.0, 2.35, BLUE_DARK),
        ("本项目目标", 3.35, 2.65, RGBColor(126, 34, 206)),
    ]
    for label, dx, dy, color in points:
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(dx), top + Inches(2.7 - dy), Inches(0.22), Inches(0.22))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.color.rgb = color
        add_textbox(slide, left + Inches(dx + 0.25), top + Inches(2.64 - dy), Inches(1.0), Inches(0.2), label, size=9.5, color=TEXT)

    add_card(slide, Inches(8.15), Inches(5.18), Inches(4.1), Inches(1.25),
             "结论", "现有 benchmark 已覆盖多步推理、代码、Web 与桌面交互，但仍缺少面向 bioinformatics workflow 的公平协议：同任务、同预算、同工具、同日志、同评判。", accent=RGBColor(126, 34, 206))
    add_footer(slide, 2)


def add_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "What Counts as a Fair Comparison?")

    add_textbox(slide, Inches(0.75), Inches(1.35), Inches(5.2), Inches(0.3),
                "公平比较的五个控制变量", size=18, bold=True, color=BLUE_DARK)

    pillars = [
        ("Task", "同一 workflow 目标与输入数据"),
        ("Environment", "同一软件栈、依赖与快照"),
        ("Tools", "同一工具权限、检索范围与接口"),
        ("Budget", "同一 token / 步数 / 时间预算"),
        ("Observation", "同一日志、I/O 观测与评分口径"),
    ]
    x0 = Inches(0.8)
    for i, (head, body) in enumerate(pillars):
        card = add_rect(slide, x0 + Inches(i * 2.45), Inches(1.9), Inches(2.1), Inches(2.1), WHITE, line=GRAY, radius=True)
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x0 + Inches(i * 2.45), Inches(1.9), Inches(2.1), Inches(0.22)).fill.solid()
        band = slide.shapes[-1]
        band.fill.fore_color.rgb = BLUE if i < 4 else RGBColor(126, 34, 206)
        band.line.color.rgb = band.fill.fore_color.rgb
        add_textbox(slide, x0 + Inches(i * 2.45 + 0.12), Inches(2.18), Inches(1.85), Inches(0.25), head, size=14, bold=True)
        add_textbox(slide, x0 + Inches(i * 2.45 + 0.12), Inches(2.58), Inches(1.85), Inches(0.95), body, size=11, color=MUTED)

    add_rect(slide, Inches(0.85), Inches(4.45), Inches(5.9), Inches(1.35), BLUE_LIGHT, line=BLUE, radius=True)
    add_textbox(slide, Inches(1.05), Inches(4.7), Inches(5.45), Inches(0.35),
                "Fairness Formula", size=18, bold=True, color=BLUE_DARK)
    add_textbox(slide, Inches(1.05), Inches(5.08), Inches(5.45), Inches(0.4),
                "Fair comparison = Same workflow + Same tools + Same budget + Same observation + Same judge", size=14, color=TEXT)

    add_textbox(slide, Inches(7.05), Inches(1.38), Inches(5.3), Inches(0.3),
                "为什么现有比较经常失真？", size=18, bold=True, color=RED)
    pitfalls = [
        "观测偏差：额外文件扫描 / 哈希 / tracing 可能改变 wall-clock 结果。",
        "环境偏差：冷启动、缓存命中、conda 激活与依赖检查开销不一致。",
        "提示偏差：隐藏 system prompt 或工具描述长度不同，会改变 token 与效果。",
        "评判偏差：只看最终成功率，忽略重试次数、I/O 开销和失败恢复。"
    ]
    add_bullet_box(slide, Inches(7.1), Inches(1.9), Inches(5.1), Inches(2.1), pitfalls, size=15)

    chart_data = CategoryChartData()
    chart_data.categories = ["Direct", "Agent"]
    chart_data.add_series("执行时间", (100, 72))
    chart_data.add_series("观测开销", (42, 8))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.15), Inches(4.25), Inches(4.9), Inches(2.0), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "示意：观测开销会扭曲比较"
    for series in chart.series:
        fill = series.format.fill
        fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE
    chart.series[1].format.fill.fore_color.rgb = AMBER
    add_footer(slide, 3)


def add_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Our Design: From Workflow Decomposition to Standardized Agent Execution")

    steps = [
        ("1. Workflow 拆解", "将 STAR/DESeq2、Kallisto/Sleuth、Varlociraptor 等流程拆成 atomic steps"),
        ("2. Skill 封装", "每个步骤定义目标、输入/输出、依赖关系和完成标准"),
        ("3. Tool 接口化", "把文件系统、执行器、日志观测统一为 agent 可调用接口"),
        ("4. 受控执行", "固定环境、预算与权限，实现 apples-to-apples 对比"),
        ("5. 统一评测", "记录 success、cost、latency、I/O、恢复与 reproducibility"),
    ]
    y = Inches(1.55)
    for i, (head, body) in enumerate(steps):
        left = Inches(0.85) + Inches(i * 2.45)
        box = add_rect(slide, left, y, Inches(2.15), Inches(2.55), WHITE, line=GRAY, radius=True)
        tag = add_rect(slide, left + Inches(0.12), y + Inches(0.15), Inches(0.9), Inches(0.3), BLUE_DARK, line=BLUE_DARK, radius=True)
        set_shape_text(tag, head.split(". ")[0], size=11, bold=True, color=WHITE)
        add_textbox(slide, left + Inches(0.15), y + Inches(0.55), Inches(1.8), Inches(0.35), head.split(". ", 1)[1], size=14, bold=True)
        add_textbox(slide, left + Inches(0.15), y + Inches(0.95), Inches(1.82), Inches(1.3), body, size=11, color=MUTED)
        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left + Inches(2.05), y + Inches(1.05), Inches(0.25), Inches(0.45))
            arr.fill.solid()
            arr.fill.fore_color.rgb = BLUE
            arr.line.color.rgb = BLUE

    add_card(slide, Inches(0.95), Inches(4.65), Inches(3.75), Inches(1.3),
             "Why Bioinformatics?", "长链条、强依赖、文件密集、环境敏感。workflow manager 能保证可复现，但 agent benchmark 还缺少与之配套的公平评测协议。", accent=GREEN)
    add_card(slide, Inches(4.85), Inches(4.65), Inches(3.75), Inches(1.3),
             "Why Skillization?", "skill 让每个步骤的目标、输入输出、完成条件和工具接口都显式化，便于 agent 执行、日志记录与误差归因。", accent=BLUE)
    add_card(slide, Inches(8.75), Inches(4.65), Inches(3.75), Inches(1.3),
             "Why It Matters?", "这样得到的不只是“谁跑通了”，而是“谁在同等约束下更稳定、更省、更可解释”。", accent=RGBColor(126, 34, 206))
    add_footer(slide, 4)


def add_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Evaluation Metrics and Expected Contributions")

    metrics = [
        ("Task success", "step-level / end-to-end 完成率"),
        ("Output quality", "结果完整性、约束满足、文件正确性"),
        ("Token & cost", "prompt / completion / tool / judge 成本"),
        ("Latency & I/O", "wall-clock、进程冷启动、I/O 扫描开销"),
        ("Recovery", "失败诊断、重试次数、回滚能力"),
        ("Reproducibility", "同环境复现、一致日志与可审计报告"),
    ]
    for i, (head, body) in enumerate(metrics):
        row = i // 3
        col = i % 3
        add_card(slide, Inches(0.75 + col * 4.05), Inches(1.45 + row * 1.35), Inches(3.7), Inches(1.1), head, body, accent=[BLUE, GREEN, AMBER, RED, BLUE_DARK, RGBColor(126, 34, 206)][i])

    chart_data = CategoryChartData()
    chart_data.categories = ["success", "quality", "cost", "latency", "recovery", "reproducibility"]
    chart_data.add_series("需要同时报告", (95, 90, 85, 80, 75, 88))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR_FILLED, Inches(0.95), Inches(4.25), Inches(4.7), Inches(2.1), chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "单一成功率不足以支撑公平评测"
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE_LIGHT
    chart.series[0].format.line.color.rgb = BLUE

    add_rect(slide, Inches(6.15), Inches(4.15), Inches(6.0), Inches(2.25), WHITE, line=GRAY, radius=True)
    add_textbox(slide, Inches(6.4), Inches(4.4), Inches(5.5), Inches(0.3), "本项目的预期贡献", size=19, bold=True, color=BLUE_DARK)
    add_bullet_box(
        slide, Inches(6.4), Inches(4.88), Inches(5.4), Inches(1.45),
        [
            "提出 workflow-level 的 fair comparison 协议，而不仅是 task-level 对比。",
            "把 bioinformatics workflow 拆解为 skill，使 agent 执行过程可控、可记录、可归因。",
            "同时报告 success、quality、token/cost、latency、I/O 与 recovery，形成更完整的 benchmark。"
        ],
        size=15
    )
    add_footer(slide, 5)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    add_slide_1(prs)
    add_slide_2(prs)
    add_slide_3(prs)
    add_slide_4(prs)
    add_slide_5(prs)

    prs.save(OUT)
    print(f"saved to {OUT}")


if __name__ == "__main__":
    main()
