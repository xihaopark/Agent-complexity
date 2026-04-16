from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/lab_workspace/projects/Agent-complexity/main/ppt_assets/outputs/academic-agent-fairness-v3.pptx")

W = Inches(13.333)
H = Inches(7.5)

BLUE = RGBColor(37, 99, 235)
BLUE_DARK = RGBColor(30, 64, 175)
BLUE_LIGHT = RGBColor(219, 234, 254)
PURPLE = RGBColor(126, 34, 206)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
GRAY = RGBColor(226, 232, 240)
GRAY_DARK = RGBColor(148, 163, 184)
WHITE = RGBColor(255, 255, 255)
SOFT = RGBColor(248, 250, 252)

FONT = "Noto Sans CJK SC"
FONT_REG = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
FONT_BOLD = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _fit_text(tf, size, bold=False, font=FONT):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    try:
        tf.fit_text(
            font_family=font,
            max_size=size,
            bold=bold,
            font_file=FONT_BOLD if bold else FONT_REG,
        )
    except Exception:
        pass


def add_textbox(slide, left, top, width, height, text, size=20, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, font=FONT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    _fit_text(tf, size, bold=bold, font=font)
    return box


def add_bullet_box(slide, left, top, width, height, bullets, size=16, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(1)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.alignment = PP_ALIGN.LEFT
        p.bullet = True
        p.space_after = Pt(6)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    _fit_text(tf, size, font=FONT)
    return box


def add_rect(slide, left, top, width, height, fill, line=GRAY, rounded=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    if rounded:
        shape.adjustments[0] = 0.10
    return shape


def set_shape_text(shape, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    _fit_text(tf, size, bold=bold)


def add_card(slide, left, top, width, height, title, body, accent=BLUE, title_size=13, body_size=10.5):
    add_rect(slide, left, top, width, height, WHITE, line=GRAY, rounded=True)
    bar = add_rect(slide, left, top, Inches(0.08), height, accent, line=accent)
    bar.line.color.rgb = accent
    add_textbox(slide, left + Inches(0.18), top + Inches(0.10), width - Inches(0.28), Inches(0.32),
                title, size=title_size, bold=True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.44), width - Inches(0.28), height - Inches(0.54),
                body, size=body_size, color=MUTED)


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.60), Inches(0.22), Inches(11.8), Inches(0.68), title, size=22, bold=True)
    line = add_rect(slide, Inches(0.60), Inches(0.86), Inches(12.10), Inches(0.03), BLUE, line=BLUE)
    line.fill.fore_color.rgb = BLUE
    if subtitle:
        add_textbox(slide, Inches(0.60), Inches(0.93), Inches(11.6), Inches(0.32), subtitle, size=10.5, color=MUTED)


def add_footer(slide, page, note="Agent Benchmarking for Bioinformatics Workflows"):
    add_textbox(slide, Inches(0.60), Inches(7.05), Inches(7.8), Inches(0.18), note, size=9, color=GRAY_DARK)
    add_textbox(slide, Inches(12.0), Inches(7.00), Inches(0.6), Inches(0.22), str(page), size=11, color=GRAY_DARK, align=PP_ALIGN.RIGHT)


def add_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Toward Fair Evaluation of AI Agents on Bioinformatics Workflows",
        "为什么需要从“能否完成任务”转向“能否在受控预算下公平地完成真实 workflow”",
    )

    add_textbox(slide, Inches(0.72), Inches(1.34), Inches(5.2), Inches(0.40), "研究动机", size=18, bold=True, color=BLUE_DARK)
    add_bullet_box(
        slide, Inches(0.74), Inches(1.75), Inches(5.45), Inches(2.2),
        [
            "AI agent 已从单轮问答扩展到多步规划、工具调用与环境交互。",
            "生物信息学 workflow 具有长链条、强依赖、文件密集和可复现性敏感等特点。",
            "现有 benchmark 很少同时控制同一 workflow、同一工具权限、同一预算与同一观测口径。",
        ],
        size=16,
    )

    top = Inches(1.45)
    x = Inches(6.65)
    labels = ["研究问题", "Agent 规划", "Workflow 执行", "结果评估"]
    for i, label in enumerate(labels):
        box = add_rect(slide, x + Inches(1.45 * i), top, Inches(1.20), Inches(0.62), BLUE_LIGHT, line=BLUE, rounded=True)
        set_shape_text(box, label, size=11.5, bold=True)
        if i < len(labels) - 1:
            chevron = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + Inches(1.45 * i + 1.18), top + Inches(0.18), Inches(0.22), Inches(0.24))
            chevron.fill.solid()
            chevron.fill.fore_color.rgb = BLUE
            chevron.line.color.rgb = BLUE

    add_card(slide, Inches(6.70), Inches(2.40), Inches(1.95), Inches(1.42),
             "GAIA", "466 real-world questions\nHuman 92%\nGPT-4+plugins 15%", accent=BLUE, body_size=10)
    add_card(slide, Inches(8.80), Inches(2.40), Inches(1.95), Inches(1.42),
             "OSWorld", "369 desktop tasks\nHuman 72.36%\nBest model 12.24%", accent=GREEN, body_size=10)
    add_card(slide, Inches(10.90), Inches(2.40), Inches(1.75), Inches(1.42),
             "BixBench", "53 bioinformatics scenarios\n296 questions\nmulti-step analysis", accent=AMBER, body_size=9.8)

    quote = add_rect(slide, Inches(0.78), Inches(4.45), Inches(11.95), Inches(1.48), BLUE_DARK, line=BLUE_DARK, rounded=True)
    set_shape_text(
        quote,
        "核心问题：如果不同 agent 在“同一 workflow”上使用不同环境、不同观测方式和不同预算，最终比较结果往往并不公平。",
        size=18,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 1)


def add_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Benchmark Landscape: What Existing Benchmarks Cover and What They Miss")

    table = slide.shapes.add_table(6, 4, Inches(0.65), Inches(1.32), Inches(7.15), Inches(4.55)).table
    headers = ["Benchmark", "主要能力", "优点", "缺口"]
    rows = [
        ["GAIA", "多步推理 + 工具", "真实世界问题", "非 workflow 场景，环境控制弱"],
        ["SWE-bench", "仓库级修复", "长链工程任务", "偏软件工程，不含数据流水线"],
        ["WebArena", "Web 交互", "可复现网页环境", "缺少文件密集型科学 workflow"],
        ["OSWorld", "桌面/多应用", "真实计算机环境", "领域约束弱，生信工具链缺失"],
        ["BixBench", "生信分析", "贴近 bioinformatics", "尚未强调公平的 agent-to-agent 协议"],
    ]
    col_widths = [1.2, 1.6, 1.55, 2.8]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(11.5)
        r.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r_idx, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else SOFT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            rr = p.runs[0]
            rr.font.name = FONT
            rr.font.size = Pt(10)
            rr.font.color.rgb = TEXT

    add_textbox(slide, Inches(8.05), Inches(1.40), Inches(4.2), Inches(0.36), "真实性 vs 领域适配度", size=14, bold=True)
    left, top, width, height = Inches(8.1), Inches(1.85), Inches(4.1), Inches(2.85)
    frame = add_rect(slide, left, top, width, height, WHITE, line=GRAY_DARK)
    frame.fill.background()
    add_rect(slide, left + width / 2, top, Inches(0.01), height, GRAY, line=GRAY)
    add_rect(slide, left, top + height / 2, width, Inches(0.01), GRAY, line=GRAY)
    add_textbox(slide, left + Inches(1.35), top + Inches(2.88), Inches(1.4), Inches(0.24), "真实性", size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, left - Inches(0.05), top + Inches(1.1), Inches(0.45), Inches(0.7), "领域\n适配度", size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    points = [
        ("GAIA", 0.90, 0.95, BLUE),
        ("SWE-bench", 1.80, 1.45, GREEN),
        ("WebArena", 2.55, 1.05, AMBER),
        ("OSWorld", 3.00, 1.45, RED),
        ("BixBench", 3.00, 2.20, BLUE_DARK),
        ("本项目目标", 3.25, 2.48, PURPLE),
    ]
    for label, dx, dy, color in points:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(dx), top + Inches(2.68 - dy), Inches(0.20), Inches(0.20))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = color
        add_textbox(slide, left + Inches(dx + 0.20), top + Inches(2.63 - dy), Inches(1.15), Inches(0.22), label, size=8.8, color=TEXT)

    add_card(
        slide, Inches(8.0), Inches(5.00), Inches(4.28), Inches(1.35),
        "结论",
        "现有 benchmark 已覆盖多步推理、代码、Web 与桌面交互，但仍缺少面向 bioinformatics workflow 的公平协议：同任务、同预算、同工具、同日志、同评判。",
        accent=PURPLE,
        body_size=10.2,
    )
    add_footer(slide, 2)


def add_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "What Counts as a Fair Comparison?")

    add_textbox(slide, Inches(0.72), Inches(1.34), Inches(5.2), Inches(0.36), "公平比较的五个控制变量", size=16, bold=True, color=BLUE_DARK)
    boxes = [
        ("Task", "同一 workflow 目标与输入数据", BLUE),
        ("Environment", "同一软件栈、依赖与快照", BLUE),
        ("Tools", "同一工具权限、检索范围与接口", BLUE),
        ("Budget", "同一 token / 步数 / 时间预算", BLUE),
        ("Observation", "同一日志、I/O 观测与评分口径", PURPLE),
    ]
    pos = [
        (Inches(0.78), Inches(1.75)), (Inches(3.10), Inches(1.75)),
        (Inches(0.78), Inches(3.05)), (Inches(3.10), Inches(3.05)),
        (Inches(1.94), Inches(4.35)),
    ]
    for (title, body, accent), (l, t) in zip(boxes, pos):
        add_card(slide, l, t, Inches(2.10), Inches(1.02), title, body, accent=accent, body_size=10)

    formula = add_rect(slide, Inches(0.78), Inches(5.62), Inches(4.45), Inches(0.72), BLUE_LIGHT, line=BLUE, rounded=True)
    set_shape_text(formula, "Fair comparison = same task + tools + budget + observation + judge", size=12.5, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(7.00), Inches(1.34), Inches(5.4), Inches(0.36), "为什么现有比较经常失真？", size=16, bold=True, color=RED)
    add_bullet_box(
        slide, Inches(7.02), Inches(1.72), Inches(5.15), Inches(2.45),
        [
            "观测偏差：额外文件扫描 / 哈希 / tracing 可能改变 wall-clock 结果。",
            "环境偏差：冷启动、缓存命中、conda 激活与依赖检查开销不一致。",
            "提示偏差：隐藏 system prompt 或工具描述长度不同，会改变 token 与效果。",
            "评判偏差：只看最终成功率，忽略重试次数、I/O 开销和失败恢复。",
        ],
        size=14,
    )

    data = CategoryChartData()
    data.categories = ["Direct", "Agent"]
    data.add_series("总执行时间", (100, 72))
    data.add_series("观测开销", (42, 8))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.08), Inches(4.55), Inches(4.85), Inches(1.55), data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "示意：观测开销会扭曲比较"
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = AMBER
    add_footer(slide, 3)


def add_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Our Design: Workflow Decomposition, Skillization, and Standardized Execution")

    steps = [
        ("1", "Workflow 拆解", "将 STAR/DESeq2、Kallisto/Sleuth、Varlociraptor 等流程拆成 atomic steps"),
        ("2", "Skill 封装", "为每个步骤定义目标、输入/输出、依赖关系和完成标准"),
        ("3", "Tool 接口化", "把文件系统、执行器、日志观测统一为 agent 可调用接口"),
        ("4", "受控执行", "固定环境、预算与权限，实现 apples-to-apples 对比"),
        ("5", "统一评测", "记录 success、cost、latency、I/O、恢复与 reproducibility"),
    ]
    positions = [
        (Inches(0.78), Inches(1.48)), (Inches(3.98), Inches(1.48)), (Inches(7.18), Inches(1.48)),
        (Inches(2.38), Inches(3.40)), (Inches(5.58), Inches(3.40)),
    ]
    for (tag, head, body), (l, t) in zip(steps, positions):
        add_rect(slide, l, t, Inches(2.75), Inches(1.50), WHITE, line=GRAY, rounded=True)
        tag_box = add_rect(slide, l + Inches(0.12), t + Inches(0.12), Inches(0.56), Inches(0.26), BLUE_DARK, line=BLUE_DARK, rounded=True)
        set_shape_text(tag_box, tag, size=11, bold=True, color=WHITE)
        add_textbox(slide, l + Inches(0.12), t + Inches(0.44), Inches(2.45), Inches(0.34), head, size=12, bold=True)
        add_textbox(slide, l + Inches(0.12), t + Inches(0.78), Inches(2.48), Inches(0.52), body, size=9.8, color=MUTED)

    connectors = [
        (Inches(3.56), Inches(2.05), Inches(0.26), Inches(0.18)),
        (Inches(6.76), Inches(2.05), Inches(0.26), Inches(0.18)),
        (Inches(3.72), Inches(3.02), Inches(0.22), Inches(0.26)),
        (Inches(6.92), Inches(3.02), Inches(0.22), Inches(0.26)),
    ]
    for l, t, w, h in connectors:
        arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, l, t, w, h)
        arr.fill.solid()
        arr.fill.fore_color.rgb = BLUE
        arr.line.color.rgb = BLUE

    add_card(slide, Inches(0.90), Inches(5.28), Inches(3.70), Inches(1.08),
             "Why Bioinformatics?", "长链条、强依赖、文件密集、环境敏感。workflow manager 能保证可复现，但 agent benchmark 仍缺少配套的公平评测协议。", accent=GREEN, body_size=9.8)
    add_card(slide, Inches(4.82), Inches(5.28), Inches(3.70), Inches(1.08),
             "Why Skillization?", "skill 让每个步骤的目标、输入输出、完成条件和工具接口都显式化，便于 agent 执行、日志记录与误差归因。", accent=BLUE, body_size=10.0)
    add_card(slide, Inches(8.74), Inches(5.28), Inches(3.70), Inches(1.08),
             "Why It Matters?", "这样得到的不只是“谁跑通了”，而是“谁在同等约束下更稳定、更省、更可解释”。", accent=PURPLE, body_size=10.0)
    add_footer(slide, 4)


def add_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Evaluation Metrics and Expected Contributions")

    metrics = [
        ("Task success", "step-level / end-to-end 完成率", BLUE),
        ("Output quality", "结果完整性、约束满足、文件正确性", GREEN),
        ("Token & cost", "prompt / completion / tool / judge 成本", AMBER),
        ("Latency & I/O", "wall-clock、冷启动、I/O 扫描开销", RED),
        ("Recovery", "失败诊断、重试次数、回滚能力", BLUE_DARK),
        ("Reproducibility", "同环境复现、一致日志与可审计报告", PURPLE),
    ]
    for i, (head, body, accent) in enumerate(metrics):
        row, col = divmod(i, 3)
        add_card(slide, Inches(0.72 + col * 4.04), Inches(1.40 + row * 1.18), Inches(3.72), Inches(0.92),
                 head, body, accent=accent, title_size=12.5, body_size=10)

    data = CategoryChartData()
    data.categories = ["success", "quality", "cost", "latency", "recovery", "reproducibility"]
    data.add_series("需要同时报告", (95, 90, 85, 80, 75, 88))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.88), Inches(4.10), Inches(4.85), Inches(2.05), data).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.has_major_gridlines = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "不能只看成功率"
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE

    add_rect(slide, Inches(6.02), Inches(3.95), Inches(6.10), Inches(2.35), WHITE, line=GRAY, rounded=True)
    add_textbox(slide, Inches(6.28), Inches(4.16), Inches(5.4), Inches(0.36), "本项目的预期贡献", size=16.5, bold=True, color=BLUE_DARK)
    add_bullet_box(
        slide, Inches(6.18), Inches(4.56), Inches(5.55), Inches(1.60),
        [
            "提出 workflow-level 的 fair comparison 协议，而不仅是 task-level 对比。",
            "把 bioinformatics workflow 拆解为 skill，使 agent 执行过程可控、可记录、可归因。",
            "同时报告 success、quality、token/cost、latency、I/O 与 recovery，形成更完整的 benchmark。",
        ],
        size=13.2,
    )
    add_footer(slide, 5)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    add_slide_1(prs)
    add_slide_2(prs)
    add_slide_3(prs)
    add_slide_4(prs)
    add_slide_5(prs)
    prs.save(OUT)
    print(f"saved to {OUT}")


if __name__ == "__main__":
    main()
