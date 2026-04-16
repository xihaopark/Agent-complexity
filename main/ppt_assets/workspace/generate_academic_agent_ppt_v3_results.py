import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/lab_workspace/projects/Agent-complexity/main/ppt_assets/outputs/academic-agent-fairness-v3-results.pptx")
RUN_BASE = Path("/lab_workspace/projects/Agent-complexity/main/finish/Renzo_DA_Agent/data/finish_run_comparisons")

W = Inches(13.333)
H = Inches(7.5)

BLUE = RGBColor(37, 99, 235)
BLUE_DARK = RGBColor(30, 64, 175)
BLUE_LIGHT = RGBColor(219, 234, 254)
PURPLE = RGBColor(126, 34, 206)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
GRAY = RGBColor(226, 232, 240)
GRAY_DARK = RGBColor(148, 163, 184)
WHITE = RGBColor(255, 255, 255)
SOFT = RGBColor(248, 250, 252)

FONT = "Noto Sans CJK SC"
FONT_REG = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
FONT_BOLD = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"


def _as_int(value, default=0):
    try:
        if value is None or value == "":
            return default
        return int(value)
    except (TypeError, ValueError):
        return default


def _as_float(value, default=0.0):
    try:
        if value is None or value == "":
            return default
        return float(value)
    except (TypeError, ValueError):
        return default


def _token_usage(entry):
    raw = entry.get("token_usage") or {}
    return {
        "prompt_tokens": _as_int(raw.get("prompt_tokens")),
        "completion_tokens": _as_int(raw.get("completion_tokens")),
        "total_tokens": _as_int(raw.get("total_tokens")),
        "reasoning_tokens": _as_int(raw.get("reasoning_tokens")),
        "cost": _as_float(raw.get("cost")),
    }


def _selected_step_id_from_llm_choice(entry):
    parsed = entry.get("parsed_output")
    if isinstance(parsed, dict):
        value = parsed.get("selected_step_id")
        if isinstance(value, str) and value.strip():
            return value.strip()
    response_text = str(((entry.get("response") or {}).get("text")) or "").strip()
    if not response_text:
        return ""
    candidate = response_text.splitlines()[0].strip().strip('"').strip("'")
    return candidate


def _llm_attributed_step_id(entry):
    value = entry.get("step_id")
    if isinstance(value, str) and value.strip():
        return value.strip()
    purpose = str(entry.get("purpose") or "")
    if purpose == "select_workflow_step":
        selected = _selected_step_id_from_llm_choice(entry)
        if selected:
            return selected
    context = entry.get("context")
    if isinstance(context, dict):
        for key in ("step_id", "current_step_id", "workflow_step_id"):
            value = context.get(key)
            if isinstance(value, str) and value.strip():
                return value.strip()
        workflow_ctx = context.get("workflow_run_context")
        if isinstance(workflow_ctx, dict):
            value = workflow_ctx.get("current_step_id")
            if isinstance(value, str) and value.strip():
                return value.strip()
    return ""


def _recompute_step_llm_usage(step_metrics, agent_run):
    by_step = {}
    turns = agent_run.get("turns") if isinstance(agent_run, dict) else []
    turns = turns if isinstance(turns, list) else []

    def add_usage(step_id, usage, calls=1):
        if not step_id:
            return
        item = by_step.setdefault(
            step_id,
            {
                "llm_call_count": 0,
                "prompt_tokens": 0,
                "completion_tokens": 0,
                "total_tokens": 0,
                "reasoning_tokens": 0,
                "cost": 0.0,
            },
        )
        item["llm_call_count"] += int(calls or 0)
        item["prompt_tokens"] += int(usage.get("prompt_tokens") or 0)
        item["completion_tokens"] += int(usage.get("completion_tokens") or 0)
        item["total_tokens"] += int(usage.get("total_tokens") or 0)
        item["reasoning_tokens"] += int(usage.get("reasoning_tokens") or 0)
        item["cost"] = round(float(item["cost"]) + float(usage.get("cost") or 0.0), 7)

    for turn in turns:
        if not isinstance(turn, dict):
            continue
        llm_usage = turn.get("llm_usage_delta") or {}
        by_node = llm_usage.get("by_node") if isinstance(llm_usage, dict) else {}
        by_node = by_node if isinstance(by_node, dict) else {}
        planner_usage = by_node.get("planner") if isinstance(by_node.get("planner"), dict) else None
        responder_usage = by_node.get("responder") if isinstance(by_node.get("responder"), dict) else None

        llm_events = turn.get("llm_trace_delta_tail") or []
        llm_events = llm_events if isinstance(llm_events, list) else []
        planner_events = [
            e for e in llm_events
            if isinstance(e, dict)
            and str(e.get("node") or "") == "planner"
            and str(e.get("purpose") or "") == "select_workflow_step"
        ]
        responder_events = [
            e for e in llm_events
            if isinstance(e, dict)
            and str(e.get("node") or "") == "responder"
            and str(e.get("purpose") or "") == "summary_reply"
        ]

        if planner_usage and planner_events:
            n = max(len(planner_events), 1)
            split = {
                "prompt_tokens": int(planner_usage.get("prompt_tokens") or 0) // n,
                "completion_tokens": int(planner_usage.get("completion_tokens") or 0) // n,
                "total_tokens": int(planner_usage.get("total_tokens") or 0) // n,
                "reasoning_tokens": int(planner_usage.get("reasoning_tokens") or 0) // n,
                "cost": float(planner_usage.get("cost") or 0.0) / n,
            }
            for e in planner_events:
                step_id = str(e.get("response_text") or "").strip().splitlines()[0].strip()
                add_usage(step_id, split, calls=1)

        if responder_usage and responder_events:
            n = max(len(responder_events), 1)
            split = {
                "prompt_tokens": int(responder_usage.get("prompt_tokens") or 0) // n,
                "completion_tokens": int(responder_usage.get("completion_tokens") or 0) // n,
                "total_tokens": int(responder_usage.get("total_tokens") or 0) // n,
                "reasoning_tokens": int(responder_usage.get("reasoning_tokens") or 0) // n,
                "cost": float(responder_usage.get("cost") or 0.0) / n,
            }
            for e in responder_events:
                step_id = str(e.get("step_id") or "").strip() or str(turn.get("current_step_id") or "").strip()
                add_usage(step_id, split, calls=1)

    output = []
    for step in step_metrics if isinstance(step_metrics, list) else []:
        if not isinstance(step, dict):
            continue
        step_id = str(step.get("step_id") or "")
        merged = dict(step)
        usage = by_step.get(step_id) or {}
        merged["llm_call_count"] = int(usage.get("llm_call_count") or 0)
        merged["prompt_tokens"] = int(usage.get("prompt_tokens") or 0)
        merged["completion_tokens"] = int(usage.get("completion_tokens") or 0)
        merged["total_tokens"] = int(usage.get("total_tokens") or 0)
        merged["reasoning_tokens"] = int(usage.get("reasoning_tokens") or 0)
        merged["cost"] = float(usage.get("cost") or 0.0)
        output.append(merged)
    return output


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _fit_text(tf, size, bold=False, font=FONT):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    try:
        tf.fit_text(
            font_family=font,
            max_size=size,
            bold=bold,
            font_file=FONT_BOLD if bold else FONT_REG,
        )
    except Exception:
        pass


def add_textbox(slide, left, top, width, height, text, size=20, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, font=FONT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    _fit_text(tf, size, bold=bold, font=font)
    return box


def add_bullet_box(slide, left, top, width, height, bullets, size=16, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(1)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.alignment = PP_ALIGN.LEFT
        p.bullet = True
        p.space_after = Pt(6)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    _fit_text(tf, size, font=FONT)
    return box


def add_rect(slide, left, top, width, height, fill, line=GRAY, rounded=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    if rounded:
        shape.adjustments[0] = 0.10
    return shape


def set_shape_text(shape, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    _fit_text(tf, size, bold=bold)


def add_card(slide, left, top, width, height, title, body, accent=BLUE, title_size=13, body_size=10.5):
    add_rect(slide, left, top, width, height, WHITE, line=GRAY, rounded=True)
    bar = add_rect(slide, left, top, Inches(0.08), height, accent, line=accent)
    bar.line.color.rgb = accent
    add_textbox(slide, left + Inches(0.18), top + Inches(0.10), width - Inches(0.28), Inches(0.32),
                title, size=title_size, bold=True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.44), width - Inches(0.28), height - Inches(0.54),
                body, size=body_size, color=MUTED)


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.60), Inches(0.22), Inches(11.8), Inches(0.68), title, size=22, bold=True)
    line = add_rect(slide, Inches(0.60), Inches(0.86), Inches(12.10), Inches(0.03), BLUE, line=BLUE)
    line.fill.fore_color.rgb = BLUE
    if subtitle:
        add_textbox(slide, Inches(0.60), Inches(0.93), Inches(11.6), Inches(0.32), subtitle, size=10.5, color=MUTED)


def add_footer(slide, page, note="Agent Benchmarking for Bioinformatics Workflows"):
    add_textbox(slide, Inches(0.60), Inches(7.05), Inches(7.8), Inches(0.18), note, size=9, color=GRAY_DARK)
    add_textbox(slide, Inches(12.0), Inches(7.00), Inches(0.6), Inches(0.22), str(page), size=11, color=GRAY_DARK, align=PP_ALIGN.RIGHT)


def load_result_bundle():
    all_runs = []
    for summary_path in sorted(RUN_BASE.glob("**/artifacts/comparison-summary.json")):
        summary = json.loads(summary_path.read_text(encoding="utf-8"))
        artifacts = summary_path.parent
        agent_run_path = artifacts / "agent-run.json"
        direct_run_path = artifacts / "direct-run.json"
        agent_run = json.loads(agent_run_path.read_text(encoding="utf-8")) if agent_run_path.exists() else {}
        direct_run = json.loads(direct_run_path.read_text(encoding="utf-8")) if direct_run_path.exists() else {}
        run_id = summary_path.parents[1].name
        all_runs.append({
            "run_id": run_id,
            "workflow_id": summary["workflow_id"],
            "direct": summary["runs"]["direct"],
            "agent": summary["runs"]["agent"],
            "direct_process": summary["processes"]["direct"],
            "agent_process": summary["processes"]["agent"],
            "step_metrics": _recompute_step_llm_usage(agent_run.get("step_metrics", []), agent_run),
            "direct_failure": direct_run.get("failure", {}),
            "agent_failure": agent_run.get("failure", {}),
        })

    latest_by_workflow = {}
    for run in all_runs:
        key = run["workflow_id"]
        if key not in latest_by_workflow or run["run_id"] > latest_by_workflow[key]["run_id"]:
            latest_by_workflow[key] = run

    success_latest = [r for r in latest_by_workflow.values() if r["agent"].get("status") == "success"]
    success_latest.sort(key=lambda r: r["workflow_id"])

    failed_runs = [r for r in all_runs if r["agent"].get("status") != "success" or r["direct"].get("status") != "success"]
    failed_runs.sort(key=lambda r: r["run_id"])

    latest_all = list(latest_by_workflow.values())
    latest_all.sort(key=lambda r: r["workflow_id"])
    return {
        "latest_all": latest_all,
        "success_latest": success_latest,
        "failed_runs": failed_runs,
    }


def short_workflow_name(workflow_id):
    mapping = {
        "snakemake-workflow-template-finish": "Template",
        "rna-seq-star-deseq2-finish": "RNA-seq STAR+DESeq2",
        "dna-seq-varlociraptor-finish": "DNA-seq Varlociraptor",
        "rna-seq-kallisto-sleuth-finish": "RNA-seq Kallisto+Sleuth",
        "zarp-finish": "Zarp",
    }
    return mapping.get(workflow_id, workflow_id)


def fmt_seconds(value):
    if value is None:
        return "-"
    return f"{float(value):.1f}"


def fmt_cost(value):
    if value is None:
        return "-"
    return f"${float(value):.4f}"


def failure_reason_text(run):
    direct = run.get("direct_failure") or {}
    agent = run.get("agent_failure") or {}
    direct_step = direct.get("step_id") or "-"
    agent_step = agent.get("step_id") or "-"
    direct_msg = (direct.get("error") or direct.get("summary") or "").replace("\n", " ").strip()
    agent_msg = (agent.get("error") or agent.get("summary") or "").replace("\n", " ").strip()
    if "nested workflow failed" in direct_msg:
        direct_msg = "nested workflow failed"
    if "nested workflow failed" in agent_msg:
        agent_msg = "nested workflow failed"
    return direct_step, agent_step, direct_msg[:60] or "-", agent_msg[:60] or "-"


def add_results_slide_1(prs, bundles):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "全部 Workflow 最终结果总览：Direct vs Agent",
        "按每个 workflow 的最新一次运行汇总；成功与失败同时展示，失败原因另页展开",
    )

    latest = bundles["latest_all"]
    rows = 1 + len(latest)
    cols = 10
    table = slide.shapes.add_table(rows, cols, Inches(0.42), Inches(1.38), Inches(12.50), Inches(2.65)).table
    headers = ["Workflow", "Latest Run", "Direct", "Agent", "Direct(s)", "Agent(s)", "Delta", "Tokens(calls)", "Cost", "Coverage"]
    widths = [2.05, 1.25, 0.72, 0.72, 0.95, 0.95, 0.82, 1.05, 0.95, 0.9]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(10.5)
        r.font.color.rgb = WHITE

    for r_idx, bundle in enumerate(latest, start=1):
        d = float(bundle["direct_process"]["duration_seconds"])
        a = float(bundle["agent_process"]["duration_seconds"])
        delta = (a - d) / d * 100 if d else 0.0
        values = [
            short_workflow_name(bundle["workflow_id"]),
            bundle["run_id"].split("-")[0],
            bundle["direct"]["status"],
            bundle["agent"]["status"],
            fmt_seconds(d),
            fmt_seconds(a),
            f"{delta:+.1f}%",
            f"{bundle['agent']['total_tokens']} ({bundle['agent'].get('llm_call_count', 0)})",
            fmt_cost(bundle["agent"]["cost"]),
            f"{bundle['agent']['declared_output_coverage']:.2f}",
        ]
        for c, text in enumerate(values):
            cell = table.cell(r_idx, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else SOFT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.runs[0]
            r.font.name = FONT
            r.font.size = Pt(9.2)
            if c in {2, 3}:
                r.font.color.rgb = GREEN if text == "success" else RED
            elif c == 6:
                r.font.color.rgb = GREEN if delta < 0 else RED
            else:
                r.font.color.rgb = TEXT

    add_card(
        slide, Inches(0.62), Inches(4.45), Inches(4.0), Inches(1.35),
        "总体观察 1",
        "4 个 workflow 已成功收敛：Template、STAR+DESeq2、Varlociraptor、Kallisto+Sleuth；当前仅 Zarp 仍失败。",
        accent=BLUE,
        body_size=10.5,
    )
    add_card(
        slide, Inches(4.72), Inches(4.45), Inches(4.0), Inches(1.35),
        "总体观察 2",
        "成功 workflow 中，Agent 总体 wall-clock 仍慢于 Direct，但 token/cost 保持在较窄区间（约 4k-5.5k tokens）。",
        accent=GREEN,
        body_size=10.5,
    )
    add_card(
        slide, Inches(8.82), Inches(4.45), Inches(4.0), Inches(1.35),
        "总体观察 3",
        "失败 workflow 需要单独分析根因：早期多失败于 prepare_references / trimming，最新失败集中在 Zarp 的 finish_target。",
        accent=PURPLE,
        body_size=10.3,
    )

    add_footer(slide, 6)


def add_results_slide_2(prs, bundles):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Step-level 结果：成功 Workflow 的内部步骤对比",
        "展示 4 个已成功 workflow 的 step_id、耗时与 step-level tokens；采用 2×2 网格避免遮挡",
    )

    success = bundles["success_latest"]
    positions = [
        (Inches(0.42), Inches(1.35), Inches(6.05)),
        (Inches(6.82), Inches(1.35), Inches(6.05)),
        (Inches(0.42), Inches(4.18), Inches(6.05)),
        (Inches(6.82), Inches(4.18), Inches(6.05)),
    ]
    for idx, bundle in enumerate(success):
        left, top, width = positions[idx]
        add_textbox(slide, left, top, width, Inches(0.36), short_workflow_name(bundle["workflow_id"]), size=11.2, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
        steps = sorted(bundle["step_metrics"], key=lambda x: x.get("declared_order", 999))
        max_steps = min(len(steps), 10)
        rows = 2 + max_steps
        table = slide.shapes.add_table(rows, 3, left, top + Inches(0.34), width, Inches(2.32)).table
        headers = ["Step", "Time(s)", "Tokens"]
        col_widths = [3.55, 1.05, 1.10]
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
        for c, text in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE_DARK
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.runs[0]
            r.font.name = FONT
            r.font.bold = True
            r.font.size = Pt(8.0)
            r.font.color.rgb = WHITE
        for r_idx, step in enumerate(steps[:max_steps], start=1):
            vals = [
                step.get("step_id", ""),
                fmt_seconds(step.get("duration_seconds_total")),
                str(step.get("total_tokens", 0)),
            ]
            for c, text in enumerate(vals):
                cell = table.cell(r_idx, c)
                cell.text = text
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 else SOFT
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                r = p.runs[0]
                r.font.name = FONT
                r.font.size = Pt(6.9)
                if c == 2 and text != "0":
                    r.font.color.rgb = PURPLE
                else:
                    r.font.color.rgb = TEXT
        total_time = sum(float(s.get("duration_seconds_total", 0) or 0) for s in steps)
        total_tokens = sum(int(s.get("total_tokens", 0) or 0) for s in steps)
        total_row = rows - 1
        totals = ["Total", fmt_seconds(total_time), str(total_tokens)]
        for c, text in enumerate(totals):
            cell = table.cell(total_row, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE_LIGHT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.runs[0]
            r.font.name = FONT
            r.font.size = Pt(7.2)
            r.font.bold = True
            r.font.color.rgb = TEXT

    add_textbox(
        slide,
        Inches(0.62), Inches(6.82), Inches(12.0), Inches(0.22),
        "观察：成功 workflow 的 token 主要集中在后段的汇总/报告相关步骤；step-level wall-clock 则更受 workflow 本身工具链长度影响。",
        size=9.2, color=MUTED, align=PP_ALIGN.LEFT
    )
    add_footer(slide, 7)


def add_results_slide_3_failures(prs, bundles):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "失败 Workflow / 失败运行：原因单独说明",
        "列出失败步骤与主要报错；区分“后续已修复成功”和“当前仍未成功”的情况",
    )

    failed = bundles["failed_runs"]
    rows = 1 + len(failed)
    cols = 7
    table = slide.shapes.add_table(rows, cols, Inches(0.38), Inches(1.45), Inches(12.55), Inches(3.4)).table
    headers = ["Run", "Workflow", "Current State", "Direct Fail Step", "Agent Fail Step", "Direct Cause", "Agent Cause"]
    widths = [1.35, 2.15, 1.2, 1.45, 1.45, 2.45, 2.45]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(9.6)
        r.font.color.rgb = WHITE
    latest_status = {r["workflow_id"]: r for r in bundles["latest_all"]}
    for r_idx, run in enumerate(failed, start=1):
        direct_step, agent_step, direct_msg, agent_msg = failure_reason_text(run)
        current_latest = latest_status[run["workflow_id"]]
        current_state = "已修复成功" if current_latest["agent"]["status"] == "success" else "仍失败"
        values = [
            run["run_id"].split("-")[0],
            short_workflow_name(run["workflow_id"]),
            current_state,
            direct_step,
            agent_step,
            direct_msg,
            agent_msg,
        ]
        for c, text in enumerate(values):
            cell = table.cell(r_idx, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else SOFT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c in {1, 5, 6} else PP_ALIGN.CENTER
            r = p.runs[0]
            r.font.name = FONT
            r.font.size = Pt(8.2)
            if c == 2:
                r.font.color.rgb = GREEN if text == "已修复成功" else RED
            else:
                r.font.color.rgb = TEXT

    add_card(
        slide, Inches(0.58), Inches(5.25), Inches(4.0), Inches(1.15),
        "失败类型 A：前置准备阶段",
        "Varlociraptor / Kallisto 早期失败均集中在 prepare_references，属于参考准备阶段的 nested workflow failure。",
        accent=AMBER,
        body_size=10.2,
    )
    add_card(
        slide, Inches(4.72), Inches(5.25), Inches(4.0), Inches(1.15),
        "失败类型 B：中间处理阶段",
        "Zarp 首次失败在 trimming，说明 read preprocessing 仍是该 workflow 的脆弱点。",
        accent=RED,
        body_size=10.2,
    )
    add_card(
        slide, Inches(8.86), Inches(5.25), Inches(4.0), Inches(1.15),
        "失败类型 C：最终交付阶段",
        "Zarp 最新失败推进到 finish_target，说明前段已打通，但 final packaging / report aggregation 仍未稳定。",
        accent=PURPLE,
        body_size=10.0,
    )
    add_footer(slide, 8)


def add_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Toward Fair Evaluation of AI Agents on Bioinformatics Workflows",
        "为什么需要从“能否完成任务”转向“能否在受控预算下公平地完成真实 workflow”",
    )

    add_textbox(slide, Inches(0.72), Inches(1.34), Inches(5.2), Inches(0.40), "研究动机", size=18, bold=True, color=BLUE_DARK)
    add_bullet_box(
        slide, Inches(0.74), Inches(1.75), Inches(5.45), Inches(2.2),
        [
            "AI agent 已从单轮问答扩展到多步规划、工具调用与环境交互。",
            "生物信息学 workflow 具有长链条、强依赖、文件密集和可复现性敏感等特点。",
            "现有 benchmark 很少同时控制同一 workflow、同一工具权限、同一预算与同一观测口径。",
        ],
        size=16,
    )

    top = Inches(1.45)
    x = Inches(6.65)
    labels = ["研究问题", "Agent 规划", "Workflow 执行", "结果评估"]
    for i, label in enumerate(labels):
        box = add_rect(slide, x + Inches(1.45 * i), top, Inches(1.20), Inches(0.62), BLUE_LIGHT, line=BLUE, rounded=True)
        set_shape_text(box, label, size=11.5, bold=True)
        if i < len(labels) - 1:
            chevron = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + Inches(1.45 * i + 1.18), top + Inches(0.18), Inches(0.22), Inches(0.24))
            chevron.fill.solid()
            chevron.fill.fore_color.rgb = BLUE
            chevron.line.color.rgb = BLUE

    add_card(slide, Inches(6.70), Inches(2.40), Inches(1.95), Inches(1.42),
             "GAIA", "466 real-world questions\nHuman 92%\nGPT-4+plugins 15%", accent=BLUE, body_size=10)
    add_card(slide, Inches(8.80), Inches(2.40), Inches(1.95), Inches(1.42),
             "OSWorld", "369 desktop tasks\nHuman 72.36%\nBest model 12.24%", accent=GREEN, body_size=10)
    add_card(slide, Inches(10.90), Inches(2.40), Inches(1.75), Inches(1.42),
             "BixBench", "53 bioinformatics scenarios\n296 questions\nmulti-step analysis", accent=AMBER, body_size=9.8)

    quote = add_rect(slide, Inches(0.78), Inches(4.45), Inches(11.95), Inches(1.48), BLUE_DARK, line=BLUE_DARK, rounded=True)
    set_shape_text(
        quote,
        "核心问题：如果不同 agent 在“同一 workflow”上使用不同环境、不同观测方式和不同预算，最终比较结果往往并不公平。",
        size=18,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 1)


def add_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Benchmark Landscape: What Existing Benchmarks Cover and What They Miss")

    table = slide.shapes.add_table(6, 4, Inches(0.65), Inches(1.32), Inches(7.15), Inches(4.55)).table
    headers = ["Benchmark", "主要能力", "优点", "缺口"]
    rows = [
        ["GAIA", "多步推理 + 工具", "真实世界问题", "非 workflow 场景，环境控制弱"],
        ["SWE-bench", "仓库级修复", "长链工程任务", "偏软件工程，不含数据流水线"],
        ["WebArena", "Web 交互", "可复现网页环境", "缺少文件密集型科学 workflow"],
        ["OSWorld", "桌面/多应用", "真实计算机环境", "领域约束弱，生信工具链缺失"],
        ["BixBench", "生信分析", "贴近 bioinformatics", "尚未强调公平的 agent-to-agent 协议"],
    ]
    col_widths = [1.2, 1.6, 1.55, 2.8]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(11.5)
        r.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r_idx, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else SOFT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            rr = p.runs[0]
            rr.font.name = FONT
            rr.font.size = Pt(10)
            rr.font.color.rgb = TEXT

    add_textbox(slide, Inches(8.05), Inches(1.40), Inches(4.2), Inches(0.36), "真实性 vs 领域适配度", size=14, bold=True)
    left, top, width, height = Inches(8.1), Inches(1.85), Inches(4.1), Inches(2.85)
    frame = add_rect(slide, left, top, width, height, WHITE, line=GRAY_DARK)
    frame.fill.background()
    add_rect(slide, left + width / 2, top, Inches(0.01), height, GRAY, line=GRAY)
    add_rect(slide, left, top + height / 2, width, Inches(0.01), GRAY, line=GRAY)
    add_textbox(slide, left + Inches(1.35), top + Inches(2.88), Inches(1.4), Inches(0.24), "真实性", size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, left - Inches(0.05), top + Inches(1.1), Inches(0.45), Inches(0.7), "领域\n适配度", size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    points = [
        ("GAIA", 0.90, 0.95, BLUE),
        ("SWE-bench", 1.80, 1.45, GREEN),
        ("WebArena", 2.55, 1.05, AMBER),
        ("OSWorld", 3.00, 1.45, RED),
        ("BixBench", 3.00, 2.20, BLUE_DARK),
        ("本项目目标", 3.25, 2.48, PURPLE),
    ]
    for label, dx, dy, color in points:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(dx), top + Inches(2.68 - dy), Inches(0.20), Inches(0.20))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = color
        add_textbox(slide, left + Inches(dx + 0.20), top + Inches(2.63 - dy), Inches(1.15), Inches(0.22), label, size=8.8, color=TEXT)

    add_card(
        slide, Inches(8.0), Inches(5.00), Inches(4.28), Inches(1.35),
        "结论",
        "现有 benchmark 已覆盖多步推理、代码、Web 与桌面交互，但仍缺少面向 bioinformatics workflow 的公平协议：同任务、同预算、同工具、同日志、同评判。",
        accent=PURPLE,
        body_size=10.2,
    )
    add_footer(slide, 2)


def add_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "What Counts as a Fair Comparison?")

    add_textbox(slide, Inches(0.72), Inches(1.34), Inches(5.2), Inches(0.36), "公平比较的五个控制变量", size=16, bold=True, color=BLUE_DARK)
    boxes = [
        ("Task", "同一 workflow 目标与输入数据", BLUE),
        ("Environment", "同一软件栈、依赖与快照", BLUE),
        ("Tools", "同一工具权限、检索范围与接口", BLUE),
        ("Budget", "同一 token / 步数 / 时间预算", BLUE),
        ("Observation", "同一日志、I/O 观测与评分口径", PURPLE),
    ]
    pos = [
        (Inches(0.78), Inches(1.75)), (Inches(3.10), Inches(1.75)),
        (Inches(0.78), Inches(3.05)), (Inches(3.10), Inches(3.05)),
        (Inches(1.94), Inches(4.35)),
    ]
    for (title, body, accent), (l, t) in zip(boxes, pos):
        add_card(slide, l, t, Inches(2.10), Inches(1.02), title, body, accent=accent, body_size=10)

    formula = add_rect(slide, Inches(0.78), Inches(5.62), Inches(4.45), Inches(0.72), BLUE_LIGHT, line=BLUE, rounded=True)
    set_shape_text(formula, "Fair comparison = same task + tools + budget + observation + judge", size=12.5, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(7.00), Inches(1.34), Inches(5.4), Inches(0.36), "为什么现有比较经常失真？", size=16, bold=True, color=RED)
    add_bullet_box(
        slide, Inches(7.02), Inches(1.72), Inches(5.15), Inches(2.45),
        [
            "观测偏差：额外文件扫描 / 哈希 / tracing 可能改变 wall-clock 结果。",
            "环境偏差：冷启动、缓存命中、conda 激活与依赖检查开销不一致。",
            "提示偏差：隐藏 system prompt 或工具描述长度不同，会改变 token 与效果。",
            "评判偏差：只看最终成功率，忽略重试次数、I/O 开销和失败恢复。",
        ],
        size=14,
    )

    data = CategoryChartData()
    data.categories = ["Direct", "Agent"]
    data.add_series("总执行时间", (100, 72))
    data.add_series("观测开销", (42, 8))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.08), Inches(4.55), Inches(4.85), Inches(1.55), data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "示意：观测开销会扭曲比较"
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = AMBER
    add_footer(slide, 3)


def add_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Our Design: Workflow Decomposition, Skillization, and Standardized Execution")

    steps = [
        ("1", "Workflow 拆解", "将 STAR/DESeq2、Kallisto/Sleuth、Varlociraptor 等流程拆成 atomic steps"),
        ("2", "Skill 封装", "为每个步骤定义目标、输入/输出、依赖关系和完成标准"),
        ("3", "Tool 接口化", "把文件系统、执行器、日志观测统一为 agent 可调用接口"),
        ("4", "受控执行", "固定环境、预算与权限，实现 apples-to-apples 对比"),
        ("5", "统一评测", "记录 success、cost、latency、I/O、恢复与 reproducibility"),
    ]
    positions = [
        (Inches(0.78), Inches(1.48)), (Inches(3.98), Inches(1.48)), (Inches(7.18), Inches(1.48)),
        (Inches(2.38), Inches(3.40)), (Inches(5.58), Inches(3.40)),
    ]
    for (tag, head, body), (l, t) in zip(steps, positions):
        add_rect(slide, l, t, Inches(2.75), Inches(1.50), WHITE, line=GRAY, rounded=True)
        tag_box = add_rect(slide, l + Inches(0.12), t + Inches(0.12), Inches(0.56), Inches(0.26), BLUE_DARK, line=BLUE_DARK, rounded=True)
        set_shape_text(tag_box, tag, size=11, bold=True, color=WHITE)
        add_textbox(slide, l + Inches(0.12), t + Inches(0.44), Inches(2.45), Inches(0.34), head, size=12, bold=True)
        add_textbox(slide, l + Inches(0.12), t + Inches(0.78), Inches(2.48), Inches(0.52), body, size=9.8, color=MUTED)

    connectors = [
        (Inches(3.56), Inches(2.05), Inches(0.26), Inches(0.18)),
        (Inches(6.76), Inches(2.05), Inches(0.26), Inches(0.18)),
        (Inches(3.72), Inches(3.02), Inches(0.22), Inches(0.26)),
        (Inches(6.92), Inches(3.02), Inches(0.22), Inches(0.26)),
    ]
    for l, t, w, h in connectors:
        arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, l, t, w, h)
        arr.fill.solid()
        arr.fill.fore_color.rgb = BLUE
        arr.line.color.rgb = BLUE

    add_card(slide, Inches(0.90), Inches(5.28), Inches(3.70), Inches(1.08),
             "Why Bioinformatics?", "长链条、强依赖、文件密集、环境敏感。workflow manager 能保证可复现，但 agent benchmark 仍缺少配套的公平评测协议。", accent=GREEN, body_size=9.8)
    add_card(slide, Inches(4.82), Inches(5.28), Inches(3.70), Inches(1.08),
             "Why Skillization?", "skill 让每个步骤的目标、输入输出、完成条件和工具接口都显式化，便于 agent 执行、日志记录与误差归因。", accent=BLUE, body_size=10.0)
    add_card(slide, Inches(8.74), Inches(5.28), Inches(3.70), Inches(1.08),
             "Why It Matters?", "这样得到的不只是“谁跑通了”，而是“谁在同等约束下更稳定、更省、更可解释”。", accent=PURPLE, body_size=10.0)
    add_footer(slide, 4)


def add_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Evaluation Metrics and Expected Contributions")

    metrics = [
        ("Task success", "step-level / end-to-end 完成率", BLUE),
        ("Output quality", "结果完整性、约束满足、文件正确性", GREEN),
        ("Token & cost", "prompt / completion / tool / judge 成本", AMBER),
        ("Latency & I/O", "wall-clock、冷启动、I/O 扫描开销", RED),
        ("Recovery", "失败诊断、重试次数、回滚能力", BLUE_DARK),
        ("Reproducibility", "同环境复现、一致日志与可审计报告", PURPLE),
    ]
    for i, (head, body, accent) in enumerate(metrics):
        row, col = divmod(i, 3)
        add_card(slide, Inches(0.72 + col * 4.04), Inches(1.40 + row * 1.18), Inches(3.72), Inches(0.92),
                 head, body, accent=accent, title_size=12.5, body_size=10)

    data = CategoryChartData()
    data.categories = ["success", "quality", "cost", "latency", "recovery", "reproducibility"]
    data.add_series("需要同时报告", (95, 90, 85, 80, 75, 88))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.88), Inches(4.10), Inches(4.85), Inches(2.05), data).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.has_major_gridlines = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "不能只看成功率"
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE

    add_rect(slide, Inches(6.02), Inches(3.95), Inches(6.10), Inches(2.35), WHITE, line=GRAY, rounded=True)
    add_textbox(slide, Inches(6.28), Inches(4.16), Inches(5.4), Inches(0.36), "本项目的预期贡献", size=16.5, bold=True, color=BLUE_DARK)
    add_bullet_box(
        slide, Inches(6.18), Inches(4.56), Inches(5.55), Inches(1.60),
        [
            "提出 workflow-level 的 fair comparison 协议，而不仅是 task-level 对比。",
            "把 bioinformatics workflow 拆解为 skill，使 agent 执行过程可控、可记录、可归因。",
            "同时报告 success、quality、token/cost、latency、I/O 与 recovery，形成更完整的 benchmark。",
        ],
        size=13.2,
    )
    add_footer(slide, 5)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    bundles = load_result_bundle()
    add_slide_1(prs)
    add_slide_2(prs)
    add_slide_3(prs)
    add_slide_4(prs)
    add_slide_5(prs)
    add_results_slide_1(prs, bundles)
    add_results_slide_2(prs, bundles)
    add_results_slide_3_failures(prs, bundles)
    prs.save(OUT)
    print(f"saved to {OUT}")




if __name__ == "__main__":
    main()
