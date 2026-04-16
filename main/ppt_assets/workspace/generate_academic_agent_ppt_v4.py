from __future__ import annotations

import json
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/lab_workspace/projects/Agent-complexity/main")
RUN_ROOT = ROOT / "finish" / "Renzo_DA_Agent" / "data" / "finish_run_comparisons"
OUT = ROOT / "ppt_assets" / "outputs" / "academic-agent-fairness-v4.pptx"

W = Inches(13.333)
H = Inches(7.5)

# Less “AI-ish” palette: near-monochrome + single muted accent
BG = RGBColor(252, 252, 253)  # off-white
TEXT = RGBColor(17, 24, 39)   # slate-900
MUTED = RGBColor(75, 85, 99)  # gray-600
LINE = RGBColor(229, 231, 235)  # gray-200
LINE_DARK = RGBColor(156, 163, 175)  # gray-400
ACCENT = RGBColor(30, 64, 175)  # indigo-800 (single accent)
ACCENT_SOFT = RGBColor(238, 242, 255)  # indigo-50
WHITE = RGBColor(255, 255, 255)

FONT = "Noto Sans CJK SC"
FONT_REG = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
FONT_BOLD = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"


def _fit_text(tf, size, bold=False, font=FONT):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    try:
        tf.fit_text(
            font_family=font,
            max_size=size,
            bold=bold,
            font_file=FONT_BOLD if bold else FONT_REG,
        )
    except Exception:
        pass


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=16, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, font=FONT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    _fit_text(tf, size, bold=bold, font=font)
    return box


def add_bullets(slide, left, top, width, height, bullets, size=12, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(1)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.alignment = PP_ALIGN.LEFT
        p.bullet = True
        p.space_after = Pt(5)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    _fit_text(tf, size, font=FONT)
    return box


def add_rect(slide, left, top, width, height, fill, line=LINE, rounded=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    if rounded:
        shape.adjustments[0] = 0.08
    return shape


def set_shape_text(shape, text, size=12, bold=False, color=TEXT, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    _fit_text(tf, size, bold=bold)


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.60), Inches(0.22), Inches(12.1), Inches(0.60), title, size=22, bold=True)
    line = add_rect(slide, Inches(0.60), Inches(0.86), Inches(12.1), Inches(0.02), LINE_DARK, line=LINE_DARK)
    line.fill.fore_color.rgb = LINE_DARK
    if subtitle:
        add_textbox(slide, Inches(0.60), Inches(0.92), Inches(12.0), Inches(0.30), subtitle, size=10.5, color=MUTED)


def add_footer(slide, page, note="Fair Agent Evaluation on Bioinformatics Workflows"):
    add_textbox(slide, Inches(0.60), Inches(7.05), Inches(8.5), Inches(0.18), note, size=9, color=LINE_DARK)
    add_textbox(slide, Inches(12.0), Inches(7.00), Inches(0.6), Inches(0.22), str(page), size=11, color=LINE_DARK, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, body, accent=False, title_size=12.5, body_size=10.5):
    add_rect(slide, left, top, width, height, BG, line=LINE, rounded=True)
    if accent:
        bar = add_rect(slide, left, top, Inches(0.06), height, ACCENT, line=ACCENT)
        bar.fill.fore_color.rgb = ACCENT
    add_textbox(slide, left + Inches(0.16), top + Inches(0.10), width - Inches(0.22), Inches(0.30),
                title, size=title_size, bold=True)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.42), width - Inches(0.22), height - Inches(0.52),
                body, size=body_size, color=MUTED)


@dataclass
class WorkflowRun:
    workflow_id: str
    output_dir: Path
    agent_run_json: Path
    comparison_json: Path
    started_at: str | None
    ended_at: str | None
    status: str
    step_metrics: list[dict[str, Any]]
    total_tokens: int | None
    cost: float | None


def _parse_dt(value: str | None) -> datetime | None:
    if not value:
        return None
    try:
        return datetime.fromisoformat(value.replace("Z", "+00:00"))
    except Exception:
        return None


def _duration_s(start: str | None, end: str | None) -> float | None:
    a = _parse_dt(start)
    b = _parse_dt(end)
    if not a or not b:
        return None
    return max(0.0, (b - a).total_seconds())


def load_successful_runs(limit: int = 2) -> list[WorkflowRun]:
    runs: list[WorkflowRun] = []
    for comp in RUN_ROOT.glob("*/artifacts/comparison-summary.json"):
        try:
            summary = json.loads(comp.read_text(encoding="utf-8"))
        except Exception:
            continue
        agent = summary.get("runs", {}).get("agent", {})
        if agent.get("status") != "success" or agent.get("workflow_status") != "success":
            continue
        agent_run = Path(summary["artifacts"]["agent_run_json"])
        if not agent_run.exists():
            continue
        try:
            ar = json.loads(agent_run.read_text(encoding="utf-8"))
        except Exception:
            continue
        runs.append(
            WorkflowRun(
                workflow_id=summary.get("workflow_id", ""),
                output_dir=Path(summary.get("output_dir", comp.parent.parent)),
                agent_run_json=agent_run,
                comparison_json=comp,
                started_at=ar.get("started_at"),
                ended_at=ar.get("ended_at"),
                status=ar.get("status", ""),
                step_metrics=ar.get("step_metrics", []),
                total_tokens=agent.get("total_tokens"),
                cost=agent.get("cost"),
            )
        )
    runs.sort(key=lambda r: r.output_dir.name, reverse=True)
    return runs[:limit]


def slide_1(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Toward Fair Evaluation of AI Agents on Bioinformatics Workflows",
        "从“能否跑通”到“在受控预算下公平对比”：workflow-level evaluation 的必要性",
    )

    add_textbox(slide, Inches(0.72), Inches(1.32), Inches(5.4), Inches(0.34),
                "Motivation", size=14, bold=True, color=ACCENT)
    add_bullets(
        slide, Inches(0.72), Inches(1.70), Inches(5.8), Inches(2.15),
        [
            "Bioinformatics workflow：多步骤、强依赖、文件密集、环境敏感。",
            "Agent benchmark 常忽略环境/观测口径差异，导致 apples-to-oranges。",
            "目标：定义可复现、可解释、可量化的 fair comparison 协议。",
        ],
        size=12,
    )

    add_textbox(slide, Inches(6.75), Inches(1.32), Inches(5.7), Inches(0.34),
                "Representative Benchmarks (signal of the gap)", size=14, bold=True, color=ACCENT)
    add_card(slide, Inches(6.75), Inches(1.70), Inches(2.0), Inches(1.05), "GAIA", "466 Qs\nHuman 92%\nGPT-4+tools 15%", accent=True, body_size=10)
    add_card(slide, Inches(8.95), Inches(1.70), Inches(2.0), Inches(1.05), "OSWorld", "369 tasks\nHuman 72.36%\nBest 12.24%", accent=False, body_size=10)
    add_card(slide, Inches(11.15), Inches(1.70), Inches(1.95), Inches(1.05), "BixBench", "Bioinformatics\n53 scenarios / 296 Qs", accent=False, body_size=10)

    quote = add_rect(slide, Inches(0.72), Inches(4.25), Inches(12.0), Inches(1.65), ACCENT_SOFT, line=LINE, rounded=True)
    set_shape_text(
        quote,
        "Key point: “Fast / Better” 结论往往被观测开销、冷启动与工具权限差异所污染。\nFairness 首先是实验协议，而不是分数。",
        size=16,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 1)


def slide_2(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Benchmark Landscape: Coverage vs. What We Need for Workflows")

    table = slide.shapes.add_table(6, 4, Inches(0.65), Inches(1.28), Inches(7.25), Inches(4.60)).table
    headers = ["Benchmark", "Focus", "Strength", "Gap for Workflow Fairness"]
    rows = [
        ["GAIA", "tool-use QA", "realistic questions", "weak environment control"],
        ["SWE-bench", "repo-level coding", "fail-to-pass tests", "not data workflow"],
        ["WebArena", "web interaction", "reproducible websites", "not file-heavy"],
        ["OSWorld", "desktop tasks", "real OS env", "domain-agnostic"],
        ["BixBench", "bioinformatics", "multi-step analysis", "fair protocol not central"],
    ]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(17, 24, 39)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(11)
        r.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r_idx, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG if r_idx % 2 else RGBColor(248, 250, 252)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            rr = p.runs[0]
            rr.font.name = FONT
            rr.font.size = Pt(10)
            rr.font.color.rgb = TEXT

    add_card(
        slide, Inches(8.10), Inches(1.28), Inches(4.40), Inches(1.35),
        "Takeaway",
        "Workflow 评测需要：同输入、同工具权限、同预算、同观测口径、同评判。\n否则无法解释差异来自 agent 还是来自实验设置。",
        accent=True,
        body_size=11,
    )
    add_card(
        slide, Inches(8.10), Inches(2.82), Inches(4.40), Inches(3.06),
        "What to report (minimum)",
        "success • quality • token/cost • latency • I/O • recovery • reproducibility",
        accent=False,
        body_size=12,
    )
    add_footer(slide, 2)


def slide_3(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Fair Comparison Protocol (Workflow-level)")

    add_textbox(slide, Inches(0.72), Inches(1.30), Inches(6.2), Inches(0.38),
                "Control Variables", size=13, bold=True, color=ACCENT)
    controls = [
        ("Task", "same workflow + same inputs"),
        ("Env", "same snapshot / deps"),
        ("Tools", "same permissions + APIs"),
        ("Budget", "same token/time/steps"),
        ("Obs", "same logging & scoring"),
    ]
    x0, y0 = Inches(0.72), Inches(1.65)
    for i, (k, v) in enumerate(controls):
        left = x0 + Inches((i % 2) * 3.15)
        top = y0 + Inches((i // 2) * 1.05)
        add_card(slide, left, top, Inches(3.00), Inches(0.90), k, v, accent=(k in {"Task", "Obs"}), body_size=11)

    add_textbox(slide, Inches(7.00), Inches(1.30), Inches(5.6), Inches(0.38),
                "Why comparisons get distorted", size=13, bold=True, color=ACCENT)
    add_bullets(
        slide, Inches(7.00), Inches(1.65), Inches(5.65), Inches(2.1),
        [
            "I/O 观测：全量扫描/哈希会改变 wall-clock。",
            "冷启动：conda 激活与依赖检查的固定开销。",
            "提示/工具描述差异：影响 token 与行为。",
            "只报 success：忽略重试与恢复成本。",
        ],
        size=12,
    )

    data = CategoryChartData()
    data.categories = ["Direct", "Agent"]
    data.add_series("Total time", (100, 72))
    data.add_series("Observation overhead", (42, 8))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.05), Inches(4.20), Inches(5.55), Inches(2.10), data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Illustration: measurement can change the outcome"
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = RGBColor(55, 65, 81)
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = RGBColor(156, 163, 175)
    add_footer(slide, 3)


def slide_4(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Our Approach: Decompose Workflow → Skillize Steps → Standardize Execution")

    steps = [
        ("1", "Decompose", "workflow → atomic steps"),
        ("2", "Skillize", "step IO contract + guardrails"),
        ("3", "Toolize", "agent-callable execution interface"),
        ("4", "Control", "same env/tools/budget/obs"),
        ("5", "Evaluate", "metrics + logs + artifacts"),
    ]
    positions = [
        (Inches(0.72), Inches(1.40)), (Inches(3.42), Inches(1.40)), (Inches(6.12), Inches(1.40)),
        (Inches(2.07), Inches(3.15)), (Inches(4.77), Inches(3.15)),
    ]
    for (tag, head, body), (l, t) in zip(steps, positions):
        add_rect(slide, l, t, Inches(2.55), Inches(1.25), BG, line=LINE, rounded=True)
        tag_box = add_rect(slide, l + Inches(0.12), t + Inches(0.12), Inches(0.46), Inches(0.24), RGBColor(17, 24, 39), line=RGBColor(17, 24, 39), rounded=True)
        set_shape_text(tag_box, tag, size=10.5, bold=True, color=WHITE)
        add_textbox(slide, l + Inches(0.12), t + Inches(0.40), Inches(2.35), Inches(0.30), head, size=12, bold=True)
        add_textbox(slide, l + Inches(0.12), t + Inches(0.70), Inches(2.40), Inches(0.42), body, size=10.5, color=MUTED)

    for l, t in [(Inches(3.25), Inches(1.95)), (Inches(5.95), Inches(1.95)), (Inches(3.58), Inches(2.85)), (Inches(6.28), Inches(2.85))]:
        chevron = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, l, t, Inches(0.22), Inches(0.18))
        chevron.fill.solid()
        chevron.fill.fore_color.rgb = LINE_DARK
        chevron.line.color.rgb = LINE_DARK

    add_card(slide, Inches(0.72), Inches(4.75), Inches(3.95), Inches(1.35), "Why workflows?",
             "Bioinformatics pipelines amplify small protocol differences.\n标准化 step 与日志是“公平评测”的前提。", accent=True, body_size=11)
    add_card(slide, Inches(4.90), Inches(4.75), Inches(3.95), Inches(1.35), "Why skills?",
             "skills make IO contracts explicit.\nEasier to attribute failures and cost to steps.", accent=False, body_size=11)
    add_card(slide, Inches(9.08), Inches(4.75), Inches(3.65), Inches(1.35), "Deliverable",
             "A workflow-level benchmark harness:\nsuccess–cost–latency with provenance.", accent=False, body_size=11)
    add_footer(slide, 4)


def slide_5(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Metrics: Report More Than Success")

    metrics = [
        ("Success", "step/end-to-end", True),
        ("Quality", "artifact completeness", False),
        ("Cost", "tokens + tools + judge", False),
        ("Latency", "wall-clock + cold start", False),
        ("I/O", "scan/hash/bytes", False),
        ("Recovery", "retries + diagnosis", True),
    ]
    for i, (head, body, accent) in enumerate(metrics):
        row, col = divmod(i, 3)
        add_card(
            slide,
            Inches(0.72 + col * 4.04),
            Inches(1.34 + row * 1.10),
            Inches(3.72),
            Inches(0.88),
            head,
            body,
            accent=accent,
            title_size=12.5,
            body_size=10.5,
        )

    data = CategoryChartData()
    data.categories = ["success", "quality", "cost", "latency", "I/O", "recovery"]
    data.add_series("report", (95, 90, 85, 80, 78, 82))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.86), Inches(3.90), Inches(4.95), Inches(2.30), data).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.has_major_gridlines = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Success alone is insufficient"
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = RGBColor(55, 65, 81)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(8)

    add_card(slide, Inches(6.20), Inches(3.90), Inches(6.45), Inches(2.30), "Expected contribution",
             "1) workflow-level fair protocol\n2) step skills + tool interfaces\n3) step-level cost/latency/recovery attribution",
             accent=True, body_size=12)
    add_footer(slide, 5)


def slide_6_results(prs: Presentation, runs: list[WorkflowRun]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Completed Workflow Results (Step-level Summary)")

    if not runs:
        add_textbox(slide, Inches(0.72), Inches(1.60), Inches(12.0), Inches(1.0),
                    "No successful runs found under finish_run_comparisons.", size=14, color=MUTED)
        add_footer(slide, 6)
        return

    # Small summary cards
    x = Inches(0.72)
    for i, r in enumerate(runs):
        dur = _duration_s(r.started_at, r.ended_at)
        dur_s = f"{dur:.1f}s" if dur is not None else "n/a"
        tokens = f"{r.total_tokens}" if r.total_tokens is not None else "n/a"
        cost = f"${r.cost:.4f}" if r.cost is not None else "n/a"
        body = f"status: {r.status}\nsteps: {len(r.step_metrics)}\ntime: {dur_s}\ntokens: {tokens} | cost: {cost}"
        add_card(slide, x + Inches(i * 4.20), Inches(1.28), Inches(4.00), Inches(1.05),
                 r.workflow_id, body, accent=True if i == 0 else False, body_size=10.5)

    # Table: workflow / step / status / duration / attempts / step tokens
    rows = 1 + sum(len(r.step_metrics) + 1 for r in runs)  # header + (workflow header + steps)
    cols = 6
    top = Inches(2.55)
    table_h = Inches(4.55)
    tbl = slide.shapes.add_table(rows, cols, Inches(0.65), top, Inches(12.05), table_h).table
    headers = ["Workflow", "Step", "Status", "Duration (s)", "Attempts", "Step tokens"]
    col_w = [2.60, 3.20, 1.10, 1.30, 1.10, 1.35]
    for i, w in enumerate(col_w):
        tbl.columns[i].width = Inches(w)
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(17, 24, 39)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = FONT
        run.font.bold = True
        run.font.size = Pt(10.5)
        run.font.color.rgb = WHITE

    r_idx = 1
    for wf in runs:
        # workflow separator row
        for c in range(cols):
            cell = tbl.cell(r_idx, c)
            cell.text = wf.workflow_id if c == 0 else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_SOFT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            if p.runs:
                rr = p.runs[0]
            else:
                rr = p.add_run()
            rr.font.name = FONT
            rr.font.size = Pt(9.8)
            rr.font.bold = True
            rr.font.color.rgb = TEXT
        r_idx += 1

        for step in sorted(wf.step_metrics, key=lambda s: s.get("declared_order", 9999)):
            step_id = step.get("step_id", "")
            status = step.get("execution_status", step.get("plan_status", ""))
            dur = step.get("duration_seconds_total", None)
            attempts = step.get("execution_attempt_count", None)
            tokens = step.get("total_tokens", None)

            values = [
                "",
                f"{step_id}",
                f"{status}",
                f"{dur:.1f}" if isinstance(dur, (int, float)) else "",
                f"{attempts}" if attempts is not None else "",
                f"{tokens}" if tokens is not None and tokens != 0 else "",
            ]
            for c, v in enumerate(values):
                cell = tbl.cell(r_idx, c)
                cell.text = v
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG if r_idx % 2 else RGBColor(248, 250, 252)
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if c in {0, 1} else PP_ALIGN.CENTER
                rr = p.runs[0] if p.runs else p.add_run()
                rr.font.name = FONT
                rr.font.size = Pt(9.2)
                rr.font.color.rgb = TEXT if c != 2 else (RGBColor(22, 101, 52) if status == "success" else RGBColor(153, 27, 27))
            r_idx += 1

    add_footer(slide, 6)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    runs = load_successful_runs(limit=2)

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6_results(prs, runs)

    prs.save(OUT)
    print(f"saved to {OUT}")


if __name__ == "__main__":
    main()
