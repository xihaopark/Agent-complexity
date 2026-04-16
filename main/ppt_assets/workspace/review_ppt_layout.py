from pathlib import Path
from typing import List

from PIL import ImageFont
from pptx import Presentation


import os

PPT_PATH = Path(os.environ.get("PPT_PATH", "/lab_workspace/projects/Agent-complexity/main/ppt_assets/outputs/academic-agent-fairness-v3.pptx"))
FONT_REG = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
FONT_BOLD = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"
EMU_PER_PT = 12700
EMU_PER_INCH = 914400


def load_font(size_pt: float, bold: bool):
    px = max(10, int(size_pt * 96 / 72))
    return ImageFont.truetype(FONT_BOLD if bold else FONT_REG, px)


def wrap_text(text: str, font, max_px: int) -> List[str]:
    lines = []
    for raw in (text or "").split("\n"):
        cur = ""
        for ch in raw:
            test = cur + ch
            if font.getlength(test) <= max_px or not cur:
                cur = test
            else:
                lines.append(cur)
                cur = ch
        lines.append(cur if cur else " ")
    return lines


def check_shape(shape):
    if not hasattr(shape, "text_frame"):
        return None
    tf = shape.text_frame
    if not tf.text.strip():
        return None

    width_px = max(20, int((shape.width / EMU_PER_INCH) * 96) - 8)
    height_px = max(20, int((shape.height / EMU_PER_INCH) * 96) - 6)
    need_px = 0

    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if runs:
            size_pt = max((r.font.size.pt if r.font.size else 14) for r in runs)
            bold = any(bool(r.font.bold) for r in runs)
        else:
            size_pt = 14
            bold = False
        font = load_font(size_pt, bold)
        lines = wrap_text(p.text, font, width_px)
        line_h = int(size_pt * 96 / 72 * 1.35)
        need_px += max(1, len(lines)) * line_h + 4

    overflow = need_px > int(height_px * 1.20)
    return {
        "text": tf.text.replace("\n", " / ")[:80],
        "need_px": need_px,
        "height_px": height_px,
        "overflow": overflow,
    }


def main():
    prs = Presentation(str(PPT_PATH))
    total_warn = 0
    for idx, slide in enumerate(prs.slides, start=1):
        warns = []
        for shape in slide.shapes:
            result = check_shape(shape)
            if result and result["overflow"]:
                text = result["text"].strip()
                if text in {"1", "2", "3", "4", "5"}:
                    continue
                if text == "Agent Benchmarking for Bioinformatics Workflows":
                    continue
                warns.append(result)
        status = "PASS" if not warns else f"WARN({len(warns)})"
        print(f"Slide {idx}: {status}")
        for w in warns:
            print(f"  - overflow: {w['need_px']}px > {w['height_px']}px | {w['text']}")
        total_warn += len(warns)
    print(f"TOTAL_WARNINGS={total_warn}")


if __name__ == "__main__":
    main()
